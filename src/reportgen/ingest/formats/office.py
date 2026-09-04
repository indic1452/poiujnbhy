"""Офисные форматы Microsoft: презентации, книги Excel, CSV и TSV.

DOCX разбирается в :mod:`reportgen.ingest.convert`; здесь — всё остальное, что
приносят с письмами:

* **PPTX / PPTM / POTX** — презентации. В технической презентации ценность
  распределена неравномерно: на слайде три слова и картинка, а под слайдом, в
  заметках докладчика, лежит объяснение, почему выбран именно этот запас по
  уровню. Поэтому заметки извлекаются наравне с телом слайда, а таблицы и
  подписи к рисункам сохраняются.
* **XLSX / XLSM / XLTX** — книги Excel. Читаются в потоковом режиме
  (``read_only``) и по **вычисленным** значениям (``data_only``): в отчёте
  нужна цифра допуска, а не строка ``=B4*1,15``. Кэша значений в файле может и
  не быть (книгу сохранил не Excel) — тогда ячейка окажется пустой, и об этом
  говорится в предупреждении.
* **XLS** — двоичные книги Excel 97–2003, через xlrd. Формат живуч: у
  измерительных приборов прошлого десятилетия выгрузка часто только такая.
* **CSV / TSV** — выгрузки приборов и биллинга. Разделитель и кодировка
  определяются автоматически: русский Excel пишет CSV с точкой с запятой и в
  cp1251.

Общее для всех табличных форматов:

* лист (страница, слайд) размечается ``page_marker`` — без этого под цитатой
  не будет ссылки «лист 3»;
* таблица выводится одним куском Markdown без пустых строк внутри: нарезчик
  корпуса режет текст по пустым строкам и таблицу пополам не разорвёт;
* очень широкая таблица режется по столбцам с повторением строки заголовков —
  Markdown-таблица на сто столбцов нечитаема ни человеком, ни моделью;
* объём одного листа ограничен, и если строки отброшены, об этом сказано и в
  тексте, и в предупреждении: молча потерянные данные хуже отсутствующих.

Тяжёлые пакеты (python-pptx, openpyxl, xlrd) импортируются внутри функций:
импорт этого модуля не должен требовать ничего стороннего.
"""

from __future__ import annotations

import csv
import datetime as dt
import io
import posixpath
import re
import xml.etree.ElementTree as ElementTree
import zipfile
from decimal import Decimal, InvalidOperation
from pathlib import Path
from typing import Any, Dict, Iterable, List, Sequence, Tuple

from ...packages import pip_hint
from .. import registry
from ..convert import ConvertedDocument, page_marker, _clean_line, _read_text

__all__ = [
    "MAX_SHEET_ROWS",
    "MAX_TABLE_COLUMNS",
    "convert_csv",
    "convert_pptx",
    "convert_xls",
    "convert_xlsx",
]

#: Сколько строк одного листа попадает в текст. Лист на сто тысяч строк — это
#: выгрузка счётчиков, а не документ: в индексе от него один шум, поэтому берём
#: начало и честно сообщаем, сколько осталось за бортом.
MAX_SHEET_ROWS = 2000

#: Шире этого таблица режется по столбцам. 25 — предел, после которого строка
#: Markdown перестаёт помещаться в контекст осмысленным куском.
MAX_TABLE_COLUMNS = 25

#: Длинный текст в рамке заголовка слайда — это уже не заголовок, а абзац.
MAX_HEADING_CHARS = 120

#: Допуск при определении порядка чтения на слайде: рамки, стоящие в одной
#: строке, редко выровнены по пикселю. 228600 EMU — четверть дюйма.
_ROW_BAND_EMU = 228600

#: Площадь объединённого диапазона, которую имеет смысл разворачивать. Excel
#: позволяет объединить столбец целиком; разворачивать миллион ячеек незачем.
_MAX_MERGE_CELLS = 20000

_OLE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

_IMAGE_NAME_RE = re.compile(
    r"^[\w \-.()]+\.(?:png|jpe?g|gif|bmp|tiff?|emf|wmf|svg|webp)$", re.IGNORECASE
)
_CELL_REF_RE = re.compile(r"^\$?([A-Za-z]{1,3})\$?([0-9]{1,7})$")
_MERGE_REF_RE = re.compile(rb'<mergeCell[^>]*\bref="([^"]+)"')

_NS_PACKAGE_REL = "{http://schemas.openxmlformats.org/package/2006/relationships}"
_NS_DOC_REL = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


# ------------------------------------------------------------ утилиты ----

def _reason(error: BaseException) -> str:
    text = str(error).strip()
    return text or error.__class__.__name__


def _looks_like_ole(path: Path) -> bool:
    """Похож ли файл на старый двоичный контейнер OLE.

    Это же начало у книги, зашифрованной паролем: Excel кладёт зашифрованный
    поток внутрь OLE, и снаружи файл перестаёт быть zip-архивом.
    """
    try:
        with open(path, "rb") as handle:
            return handle.read(8) == _OLE_MAGIC
    except OSError:
        return False


def _package_hint(path: Path) -> str:
    if _looks_like_ole(path):
        return (
            " — внутри старый двоичный формат или файл зашифрован паролем; "
            "снимите пароль и пересохраните в современном формате"
        )
    return ""


def _number_text(value: float) -> str:
    """Число как в ячейке: без научной нотации и без хвостовых нулей.

    ``1.234e-06`` в отчёте о затухании читается как опечатка, а не как число,
    поэтому выводим ``0.000001234``.
    """
    if value != value:  # NaN
        return "не число"
    if value in (float("inf"), float("-inf")):
        return "бесконечность" if value > 0 else "минус бесконечность"
    try:
        decimal = Decimal(repr(float(value)))
    except (InvalidOperation, ValueError, OverflowError):  # pragma: no cover — защита
        return str(value)
    text = format(decimal, "f")
    if "." in text:
        text = text.rstrip("0").rstrip(".")
    return text or "0"


def _duration_text(value: dt.timedelta) -> str:
    total = int(value.total_seconds())
    sign = "-" if total < 0 else ""
    total = abs(total)
    hours, rest = divmod(total, 3600)
    minutes, seconds = divmod(rest, 60)
    return f"{sign}{hours:d}:{minutes:02d}:{seconds:02d}"


def _cell_text(value: Any) -> str:
    """Значение ячейки → строка. Даты в ISO, числа без научной нотации."""
    if value is None:
        return ""
    if isinstance(value, bool):
        return "истина" if value else "ложь"
    if isinstance(value, dt.datetime):
        if value.time() == dt.time(0, 0):
            return value.date().isoformat()
        return value.isoformat(sep=" ")
    if isinstance(value, dt.date):
        return value.isoformat()
    if isinstance(value, dt.time):
        return value.isoformat()
    if isinstance(value, dt.timedelta):
        return _duration_text(value)
    if isinstance(value, int):
        return str(value)
    if isinstance(value, float):
        return _number_text(value)
    if isinstance(value, Decimal):
        return _number_text(float(value))
    return _clean_line(str(value))


def _escape_cell(text: str) -> str:
    """Готовит текст к вставке в ячейку Markdown-таблицы."""
    text = text.replace("\r\n", "\n").replace("\r", "\n").replace("\t", " ")
    parts = [part.strip() for part in text.split("\n")]
    joined = " <br> ".join(part for part in parts if part)
    return joined.replace("|", "\\|").strip()


def _row_is_empty(row: Sequence[str]) -> bool:
    return not any(cell.strip() for cell in row)


def _trim_table(rows: List[List[str]]) -> Tuple[List[List[str]], int, int]:
    """Обрезает пустые строки и столбцы по краям.

    Возвращает таблицу и смещения (сколько строк сверху и столбцов слева
    отброшено) — они нужны, чтобы попадать в исходные координаты листа.
    """
    first_row = 0
    last_row = len(rows)
    while first_row < last_row and _row_is_empty(rows[first_row]):
        first_row += 1
    while last_row > first_row and _row_is_empty(rows[last_row - 1]):
        last_row -= 1
    body = [list(row) for row in rows[first_row:last_row]]
    if not body:
        return [], first_row, 0

    width = max(len(row) for row in body)
    body = [row + [""] * (width - len(row)) for row in body]
    first_column = 0
    last_column = width
    while first_column < last_column and not any(
        row[first_column].strip() for row in body
    ):
        first_column += 1
    while last_column > first_column and not any(
        row[last_column - 1].strip() for row in body
    ):
        last_column -= 1
    body = [row[first_column:last_column] for row in body]
    return body, first_row, first_column


def _column_groups(width: int) -> List[Tuple[int, int]]:
    """Границы кусков, на которые режется слишком широкая таблица."""
    if width <= MAX_TABLE_COLUMNS:
        return [(0, width)]
    return [
        (start, min(start + MAX_TABLE_COLUMNS, width))
        for start in range(0, width, MAX_TABLE_COLUMNS)
    ]


def _table_blocks(rows: Sequence[Sequence[str]]) -> List[str]:
    """Таблица → куски Markdown.

    Первая строка считается строкой заголовков. Если столбцов больше
    :data:`MAX_TABLE_COLUMNS`, таблица режется по столбцам, и заголовки
    повторяются в каждом куске — иначе продолжение читается как набор цифр без
    смысла. Внутри куска нет пустых строк: нарезчик корпуса режет по ним и
    таблицу пополам не разорвёт.
    """
    body = [list(row) for row in rows if row is not None]
    if not body:
        return []
    width = max((len(row) for row in body), default=0)
    if not width:
        return []
    body = [
        [_escape_cell(cell) for cell in row] + [""] * (width - len(row))
        for row in body
    ]

    groups = _column_groups(width)
    blocks: List[str] = []
    header = body[0]
    for start, stop in groups:
        if len(groups) > 1:
            blocks.append(f"### Столбцы {start + 1}–{stop}")
        lines = [
            "| " + " | ".join(header[start:stop]) + " |",
            "| " + " | ".join(["---"] * (stop - start)) + " |",
        ]
        for row in body[1:]:
            lines.append("| " + " | ".join(row[start:stop]) + " |")
        blocks.append("\n".join(lines))
    return blocks


def _expand_merges(
    grid: List[List[str]],
    ranges: Iterable[Tuple[int, int, int, int]],
    *,
    first_row: int,
    first_column: int,
) -> int:
    """Разворачивает объединённые ячейки: значение из первой ячейки диапазона
    проставляется во все ячейки этого диапазона.

    В файле значение объединённой ячейки хранится только в левой верхней
    ячейке, остальные пустые. В таблице допусков это означает шапку вида
    «Диапазон частот | (пусто) | (пусто)»: столбцы теряют название, и строка
    таблицы перестаёт быть самодостаточной. Диапазоны задаются в координатах
    листа (с единицы), ``first_row`` и ``first_column`` — координаты левой
    верхней ячейки ``grid``.
    """
    if not grid:
        return 0
    height = len(grid)
    width = max(len(row) for row in grid)
    filled = 0
    for min_column, min_row, max_column, max_row in ranges:
        if (max_row - min_row + 1) * (max_column - min_column + 1) > _MAX_MERGE_CELLS:
            continue
        top = min_row - first_row
        left = min_column - first_column
        if top < 0 or left < 0 or top >= height or left >= width:
            continue
        source = grid[top][left] if left < len(grid[top]) else ""
        if not source.strip():
            continue
        bottom = min(max_row - first_row, height - 1)
        right = min(max_column - first_column, width - 1)
        for row_index in range(top, bottom + 1):
            row = grid[row_index]
            row.extend([""] * (width - len(row)))
            for column_index in range(left, right + 1):
                if not row[column_index].strip():
                    row[column_index] = source
                    filled += 1
    return filled


def _limit_note(shown: int, skipped: int) -> str:
    return (
        f"*Показаны первые {shown} строк, пропущено строк: {skipped}. "
        f"Полные данные — в исходном файле.*"
    )


def _sheet_pieces(
    name: str,
    index: int,
    rows: Sequence[Sequence[str]],
    skipped: int,
) -> List[str]:
    """Раздел одного листа: заголовок, маркер страницы, таблица, примечание."""
    blocks = _table_blocks(rows)
    if not blocks:
        return []
    pieces = [f"## Лист «{name}»", page_marker(index)]
    pieces.extend(blocks)
    if skipped:
        pieces.append(_limit_note(len(rows), skipped))
    return pieces


# --------------------------------------------- объединённые ячейки XLSX ---

def _column_number(letters: str) -> int:
    number = 0
    for character in letters.upper():
        number = number * 26 + (ord(character) - 64)
    return number


def _range_boundaries(ref: str) -> Tuple[int, int, int, int] | None:
    """«B2:D4» → (первый столбец, первая строка, последний столбец, последняя строка)."""
    parts = ref.replace("$", "").split(":")
    if not parts or len(parts) > 2:
        return None
    corners: List[Tuple[int, int]] = []
    for part in parts:
        match = _CELL_REF_RE.match(part.strip())
        if match is None:
            return None
        corners.append((_column_number(match.group(1)), int(match.group(2))))
    if len(corners) == 1:
        corners.append(corners[0])
    (first_column, first_row), (last_column, last_row) = corners
    return (
        min(first_column, last_column),
        min(first_row, last_row),
        max(first_column, last_column),
        max(first_row, last_row),
    )


def _sheet_parts(archive: zipfile.ZipFile) -> List[Tuple[str, str]]:
    """Пары «имя листа → путь внутри архива» из xl/workbook.xml."""
    book = ElementTree.fromstring(archive.read("xl/workbook.xml"))
    targets: Dict[str, str] = {}
    try:
        relations = ElementTree.fromstring(archive.read("xl/_rels/workbook.xml.rels"))
    except KeyError:
        return []
    for node in relations:
        if node.tag == f"{_NS_PACKAGE_REL}Relationship":
            identifier = node.get("Id")
            target = node.get("Target")
            if identifier and target:
                targets[identifier] = target
    parts: List[Tuple[str, str]] = []
    for node in book.iter():
        if node.tag.rsplit("}", 1)[-1] != "sheet":
            continue
        name = node.get("name")
        target = targets.get(node.get(f"{_NS_DOC_REL}id", ""))
        if not name or not target:
            continue
        part = target.lstrip("/") if target.startswith("/") else posixpath.normpath(
            posixpath.join("xl", target)
        )
        parts.append((name, part))
    return parts


def _scan_merge_refs(archive: zipfile.ZipFile, part: str) -> List[str]:
    """Ссылки на объединённые диапазоны листа, потоковым чтением.

    XML листа читается блоками и не разбирается целиком: у выгрузки на сто
    тысяч строк дерево не поместится в память, а нужны из него четыре десятка
    коротких элементов ``mergeCell``. Совпадение внутри текста ячейки
    невозможно: там ``<`` записан как ``&lt;``.
    """
    found: Dict[str, None] = {}
    tail = b""
    with archive.open(part) as stream:
        while True:
            block = stream.read(1 << 20)
            if not block:
                break
            data = tail + block
            for match in _MERGE_REF_RE.finditer(data):
                found[match.group(1).decode("ascii", "ignore")] = None
            tail = data[-200:]
    return list(found)


def _merged_ranges(path: Path) -> Dict[str, List[Tuple[int, int, int, int]]]:
    """Объединённые диапазоны по листам книги.

    Потоковый режим openpyxl объединения не отдаёт (``ReadOnlyWorksheet`` их не
    хранит), а отказываться от потокового режима из-за шапки таблицы нельзя.
    Поэтому диапазоны берутся напрямую из архива. Любая неудача здесь не
    фатальна: без разворачивания таблица останется с пустыми ячейками в шапке,
    но текст будет извлечён.
    """
    result: Dict[str, List[Tuple[int, int, int, int]]] = {}
    try:
        with zipfile.ZipFile(path) as archive:
            for name, part in _sheet_parts(archive):
                try:
                    refs = _scan_merge_refs(archive, part)
                except (KeyError, OSError, zipfile.BadZipFile):
                    continue
                boundaries = [
                    item for item in (_range_boundaries(ref) for ref in refs)
                    if item is not None
                ]
                if boundaries:
                    result[name] = boundaries
    except Exception:  # noqa: BLE001 — вспомогательный разбор приёма не роняет
        return {}
    return result


# ----------------------------------------------------------------- PPTX ---

def _shape_alt_text(shape: Any) -> str:
    """Альтернативный текст фигуры (поле «Описание» в PowerPoint).

    Для схемы тракта, вставленной картинкой, это единственный текст, который
    вообще можно извлечь, — терять его нельзя. Автоматическое описание вида
    ``image3.png`` смысла не несёт и отбрасывается.
    """
    try:
        for node in shape._element.iter():
            if node.tag.rsplit("}", 1)[-1] != "cNvPr":
                continue
            for key in ("descr", "title"):
                value = _clean_line(node.get(key) or "")
                if value and not _IMAGE_NAME_RE.match(value):
                    return value
            break
    except Exception:  # noqa: BLE001 — экзотическая фигура не должна ронять слайд
        return ""
    return ""


def _text_frame_block(frame: Any) -> str:
    """Текст рамки: абзац — строка, вложенный уровень — пункт списка."""
    lines: List[str] = []
    for paragraph in frame.paragraphs:
        try:
            raw = paragraph.text
            level = int(getattr(paragraph, "level", 0) or 0)
        except Exception:  # noqa: BLE001 — битый абзац пропускаем
            continue
        for part in raw.replace("\v", "\n").split("\n"):
            text = _clean_line(part)
            if not text:
                continue
            if level > 0:
                lines.append("  " * (level - 1) + "- " + text)
            else:
                lines.append(text)
    return "\n".join(lines)


def _pptx_table_rows(table: Any) -> List[List[str]]:
    """Таблица слайда → строки текста, с разворачиванием объединённых ячеек."""
    rows: List[List[str]] = []
    for row in table.rows:
        cells: List[str] = []
        for cell in row.cells:
            try:
                cells.append(_clean_line(cell.text.replace("\v", " ")))
            except Exception:  # noqa: BLE001 — ячейка без текстовой рамки
                cells.append("")
        rows.append(cells)
    try:
        _expand_pptx_merges(table, rows)
    except Exception:  # noqa: BLE001 — старая версия python-pptx без сведений о слиянии
        pass
    return rows


def _expand_pptx_merges(table: Any, rows: List[List[str]]) -> None:
    for row_index, row in enumerate(table.rows):
        for column_index, cell in enumerate(row.cells):
            if not getattr(cell, "is_merge_origin", False):
                continue
            value = rows[row_index][column_index]
            if not value.strip():
                continue
            height = int(getattr(cell, "span_height", 1) or 1)
            width = int(getattr(cell, "span_width", 1) or 1)
            for offset_row in range(row_index, min(row_index + height, len(rows))):
                target = rows[offset_row]
                for offset_column in range(
                    column_index, min(column_index + width, len(target))
                ):
                    if not target[offset_column].strip():
                        target[offset_column] = value


def _chart_blocks(shape: Any) -> List[str]:
    """Диаграмма слайда: название и, если получится, ряды данных таблицей.

    Числа с графика в отчёт не переносятся (это инвариант: значения приходят из
    факт-пакета), но для поиска подпись и категории важны — по ним слайд
    находится.
    """
    title = ""
    try:
        chart = shape.chart
        if chart.has_title:
            title = _clean_line(chart.chart_title.text_frame.text)
    except Exception:  # noqa: BLE001 — диаграмма без доступного названия
        return ["**Диаграмма**"]
    blocks = [f"**Диаграмма: {title}**" if title else "**Диаграмма**"]
    try:
        plot = chart.plots[0]
        categories = [_clean_line(str(item)) for item in plot.categories]
        series = list(plot.series)
        if not categories or not series:
            return blocks
        rows: List[List[str]] = [
            ["Категория"]
            + [
                _clean_line(str(item.name or "")) or f"Ряд {number}"
                for number, item in enumerate(series, start=1)
            ]
        ]
        for index, category in enumerate(categories):
            values = [
                _cell_text(item.values[index]) if index < len(item.values) else ""
                for item in series
            ]
            rows.append([category, *values])
        blocks.extend(_table_blocks(rows))
    except Exception:  # noqa: BLE001 — экзотический тип диаграммы
        return blocks
    return blocks


def _shape_top_left(shape: Any) -> Tuple[int, int]:
    try:
        top = int(shape.top or 0)
    except Exception:  # noqa: BLE001 — размещение унаследовано от макета
        top = 0
    try:
        left = int(shape.left or 0)
    except Exception:  # noqa: BLE001 — размещение унаследовано от макета
        left = 0
    return top, left


def _ordered_shapes(shapes: Iterable[Any], skip_id: int | None = None) -> List[Any]:
    """Фигуры в порядке чтения: сверху вниз, в пределах строки — слева направо.

    Порядок хранения в файле — это порядок наложения (z-order), он совпадает с
    порядком чтения только случайно: подпись, добавленная последней, окажется в
    конце слайда, хотя читается первой.
    """
    items: List[Tuple[int, int, int, Any]] = []
    for index, shape in enumerate(shapes):
        try:
            if skip_id is not None and shape.shape_id == skip_id:
                continue
        except Exception:  # noqa: BLE001 — фигура без идентификатора
            pass
        top, left = _shape_top_left(shape)
        items.append((top // _ROW_BAND_EMU, left, index, shape))
    items.sort(key=lambda item: (item[0], item[1], item[2]))
    return [item[3] for item in items]


def _shape_blocks(shape: Any, state: Dict[str, int]) -> List[str]:
    """Один объект слайда → куски Markdown."""
    try:
        shape_type = str(shape.shape_type or "")
    except Exception:  # noqa: BLE001 — нестандартная фигура
        shape_type = ""

    if "GROUP" in shape_type:
        blocks: List[str] = []
        for child in _ordered_shapes(shape.shapes):
            blocks.extend(_shape_blocks(child, state))
        return blocks

    if getattr(shape, "has_table", False):
        return _table_blocks(_pptx_table_rows(shape.table))

    if getattr(shape, "has_chart", False):
        return _chart_blocks(shape)

    if "PICTURE" in shape_type:
        state["images"] = state.get("images", 0) + 1
        alt = _shape_alt_text(shape)
        return [f"![{alt}]()" if alt else f"![рисунок {state['images']}]()"]

    if getattr(shape, "has_text_frame", False):
        text = _text_frame_block(shape.text_frame)
        if text:
            return [text]

    alt = _shape_alt_text(shape)
    return [alt] if alt else []


def _slide_title(slide: Any) -> Tuple[str, int | None]:
    try:
        shape = slide.shapes.title
    except Exception:  # noqa: BLE001 — макет без рамки заголовка
        return "", None
    if shape is None:
        return "", None
    try:
        identifier = int(shape.shape_id)
    except Exception:  # noqa: BLE001 — фигура без идентификатора
        identifier = None
    if not getattr(shape, "has_text_frame", False):
        return "", identifier
    return _clean_line(shape.text_frame.text.replace("\v", " ")), identifier


def _notes_text(slide: Any) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        frame = slide.notes_slide.notes_text_frame
    except Exception:  # noqa: BLE001 — повреждённая страница заметок
        return ""
    if frame is None:
        return ""
    return _text_frame_block(frame)


def _slide_blocks(slide: Any, state: Dict[str, int], warnings: List[str], number: int) -> List[str]:
    blocks: List[str] = []
    title, title_id = _slide_title(slide)
    if title:
        if len(title) <= MAX_HEADING_CHARS:
            blocks.append(f"### {title}")
        else:
            blocks.append(title)

    try:
        shapes = _ordered_shapes(slide.shapes, title_id)
    except Exception as error:  # noqa: BLE001 — слайд без читаемого дерева фигур
        warnings.append(f"слайд {number}: фигуры не прочитаны ({_reason(error)})")
        shapes = []
    for shape in shapes:
        try:
            blocks.extend(_shape_blocks(shape, state))
        except Exception as error:  # noqa: BLE001 — одна фигура не роняет слайд
            warnings.append(f"слайд {number}: объект пропущен ({_reason(error)})")

    notes = _notes_text(slide)
    if notes:
        blocks.append("### Заметки к слайду")
        blocks.append(notes)
    return [block for block in blocks if block.strip()]


def convert_pptx(path: Path) -> ConvertedDocument:
    """PPTX/POTX → Markdown: слайд = раздел, заметки докладчика сохраняются."""
    result = ConvertedDocument(title=path.stem, meta={"source_format": "pptx"})
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        result.warnings.append(
            "для разбора презентаций нужен пакет python-pptx "
            "(%s); " % pip_hint("python-pptx") +
            "в изолированном контуре ставится из локального зеркала"
        )
        return result

    try:
        presentation = Presentation(str(path))
    except Exception as error:  # noqa: BLE001 — битый файл не роняет приём каталога
        result.warnings.append(
            f"не удалось открыть презентацию: {_reason(error)}{_package_hint(path)}"
        )
        return result

    state: Dict[str, int] = {"images": 0}
    pieces: List[str] = []
    empty: List[int] = []
    try:
        slides = list(presentation.slides)
    except Exception as error:  # noqa: BLE001 — повреждён список слайдов
        result.warnings.append(f"список слайдов не прочитан: {_reason(error)}")
        return result

    result.page_count = len(slides)
    for number, slide in enumerate(slides, start=1):
        try:
            blocks = _slide_blocks(slide, state, result.warnings, number)
        except Exception as error:  # noqa: BLE001 — один слайд не роняет презентацию
            result.warnings.append(f"слайд {number}: разобран не полностью ({_reason(error)})")
            blocks = []
        if not blocks:
            empty.append(number)
            continue
        pieces.append(f"## Слайд {number}")
        pieces.append(page_marker(number))
        pieces.extend(blocks)

    result.text = "\n\n".join(pieces)

    if state["images"]:
        result.meta["images"] = state["images"]
    if empty:
        result.meta["empty_slides"] = len(empty)
        listing = ", ".join(str(item) for item in empty[:20])
        if len(empty) > 20:
            listing += ", …"
        result.warnings.append(
            f"слайдов без содержимого: {len(empty)} (№ {listing}) — они пропущены, "
            "нумерация остальных сохранена"
        )
    if not slides:
        result.warnings.append("в презентации нет ни одного слайда")

    core_title, author = "", ""
    try:
        properties = presentation.core_properties
        core_title = _clean_line(properties.title or "")
        author = _clean_line(properties.author or "")
    except Exception:  # noqa: BLE001 — пакет без свойств документа
        pass
    if author:
        result.meta["author"] = author
    first_title = ""
    for slide in slides:
        first_title, _ = _slide_title(slide)
        if first_title:
            break
    result.title = core_title or first_title or path.stem

    if result.is_empty:
        result.warnings.append("в презентации не найдено текста")
    return result


# ----------------------------------------------------------------- XLSX ---

def _read_worksheet(worksheet: Any) -> Tuple[List[List[str]], int]:
    """Лист → таблица текста и число отброшенных строк.

    Потоковый ``iter_rows`` всегда добивает строки пустыми ячейками от A1, а не
    от первой заполненной ячейки, поэтому левый верхний угол полученной таблицы
    — это всегда ячейка A1 листа. На этом держится пересчёт координат
    объединённых диапазонов.
    """
    rows: List[List[str]] = []
    skipped = 0
    for values in worksheet.iter_rows(values_only=True):
        if len(rows) >= MAX_SHEET_ROWS:
            if any(value is not None and str(value).strip() for value in values):
                skipped += 1
            continue
        rows.append([_cell_text(value) for value in values])
    return rows, skipped


def convert_xlsx(path: Path) -> ConvertedDocument:
    """XLSX/XLSM/XLTX → Markdown: лист = раздел, значения — вычисленные."""
    result = ConvertedDocument(title=path.stem, meta={"source_format": "xlsx"})
    try:
        import openpyxl  # type: ignore
    except ImportError:
        result.warnings.append(
            "для разбора книг Excel нужен пакет openpyxl "
            "(%s); " % pip_hint("openpyxl") +
            "в изолированном контуре ставится из локального зеркала"
        )
        return result

    try:
        # read_only — потоковое чтение: книга на сотню листов не должна
        # разворачиваться в память целиком. data_only — вычисленные значения:
        # в отчёте нужна цифра допуска, а не строка «=B4*1,15».
        workbook = openpyxl.load_workbook(
            str(path), read_only=True, data_only=True, keep_links=False
        )
    except Exception as error:  # noqa: BLE001 — битый файл не роняет приём каталога
        result.warnings.append(
            f"не удалось открыть книгу Excel: {_reason(error)}{_package_hint(path)}"
        )
        return result

    merged = _merged_ranges(path)
    pieces: List[str] = []
    empty: List[str] = []
    hidden: List[str] = []
    try:
        names = list(workbook.sheetnames)
        result.page_count = len(names)
        for index, name in enumerate(names, start=1):
            try:
                worksheet = workbook[name]
                if str(getattr(worksheet, "sheet_state", "visible")) != "visible":
                    hidden.append(name)
                rows, skipped = _read_worksheet(worksheet)
            except Exception as error:  # noqa: BLE001 — один лист не роняет книгу
                result.warnings.append(f"лист «{name}»: не прочитан ({_reason(error)})")
                continue

            if name in merged:
                try:
                    _expand_merges(rows, merged[name], first_row=1, first_column=1)
                except Exception as error:  # noqa: BLE001 — кривые границы диапазона
                    result.warnings.append(
                        f"лист «{name}»: объединённые ячейки не развёрнуты ({_reason(error)})"
                    )

            trimmed, _, _ = _trim_table(rows)
            section = _sheet_pieces(name, index, trimmed, skipped)
            if not section:
                empty.append(name)
                continue
            pieces.extend(section)
            if skipped:
                result.warnings.append(
                    f"лист «{name}»: в текст взяты первые {len(trimmed)} строк, "
                    f"пропущено {skipped} — разбейте выгрузку на части, "
                    "если нужны все данные"
                )
    except Exception as error:  # noqa: BLE001 — книга разобрана не полностью
        result.warnings.append(f"книга Excel разобрана не полностью: {_reason(error)}")
    finally:
        try:
            workbook.close()
        except Exception:  # pragma: no cover — закрытие уже закрытой книги
            pass

    result.text = "\n\n".join(pieces)
    if empty:
        result.meta["empty_sheets"] = len(empty)
        result.warnings.append(
            "листов без данных: " + ", ".join(f"«{name}»" for name in empty[:20])
        )
    if hidden:
        result.warnings.append(
            "скрытые листы разобраны наравне с остальными: "
            + ", ".join(f"«{name}»" for name in hidden[:20])
        )

    core_title, author = "", ""
    try:
        properties = workbook.properties
        core_title = _clean_line(properties.title or "")
        author = _clean_line(properties.creator or "")
    except Exception:  # noqa: BLE001 — книга без свойств документа
        pass
    if author:
        result.meta["author"] = author
    result.title = core_title or path.stem

    if result.is_empty:
        result.warnings.append(
            "в книге Excel нет данных: все листы пусты либо в файле нет "
            "вычисленных значений формул (пересохраните книгу в Excel)"
        )
    return result


# ------------------------------------------------------------------ XLS ---

def _xls_cell_text(book: Any, sheet: Any, row: int, column: int) -> str:
    import xlrd  # type: ignore

    kind = sheet.cell_type(row, column)
    value = sheet.cell_value(row, column)
    if kind in (xlrd.XL_CELL_EMPTY, xlrd.XL_CELL_BLANK):
        return ""
    if kind == xlrd.XL_CELL_BOOLEAN:
        return "истина" if value else "ложь"
    if kind == xlrd.XL_CELL_ERROR:
        return str(xlrd.error_text_from_code.get(value, "#ОШИБКА"))
    if kind == xlrd.XL_CELL_DATE:
        try:
            moment = xlrd.xldate.xldate_as_datetime(value, book.datemode)
        except Exception:  # noqa: BLE001 — дата вне допустимого диапазона Excel
            return _cell_text(value)
        if float(value) < 1:
            return moment.time().isoformat()
        return _cell_text(moment)
    return _cell_text(value)


def _read_xls_sheet(book: Any, sheet: Any) -> Tuple[List[List[str]], int]:
    rows: List[List[str]] = []
    skipped = 0
    for row in range(sheet.nrows):
        if len(rows) >= MAX_SHEET_ROWS:
            skipped += 1
            continue
        rows.append(
            [_xls_cell_text(book, sheet, row, column) for column in range(sheet.ncols)]
        )
    return rows, skipped


def convert_xls(path: Path) -> ConvertedDocument:
    """XLS (Excel 97–2003) → Markdown. Формат читает только xlrd."""
    result = ConvertedDocument(title=path.stem, meta={"source_format": "xls"})
    try:
        import xlrd  # type: ignore
    except ImportError:
        result.warnings.append(
            "для разбора старых книг .xls нужен пакет xlrd "
            "(%s); " % pip_hint("xlrd") +
            "современные .xlsx он не читает, это отдельный конвертер"
        )
        return result

    log = io.StringIO()
    book = None
    failure = ""
    # formatting_info даёт сведения об объединённых ячейках; на части файлов он
    # не поддерживается, и тогда книга читается без него — лучше таблица с
    # пустой шапкой, чем ничего.
    for formatting in (True, False):
        try:
            book = xlrd.open_workbook(
                str(path), formatting_info=formatting, logfile=log, on_demand=False
            )
            break
        except Exception as error:  # noqa: BLE001 — битый файл не роняет приём
            failure = _reason(error)
    if book is None:
        result.warnings.append(f"не удалось открыть книгу .xls: {failure}")
        return result

    pieces: List[str] = []
    empty: List[str] = []
    try:
        sheets = book.sheets()
        result.page_count = len(sheets)
        for index, sheet in enumerate(sheets, start=1):
            name = str(sheet.name)
            try:
                rows, skipped = _read_xls_sheet(book, sheet)
            except Exception as error:  # noqa: BLE001 — один лист не роняет книгу
                result.warnings.append(f"лист «{name}»: не прочитан ({_reason(error)})")
                continue
            ranges = [
                (column_low + 1, row_low + 1, column_high, row_high)
                for row_low, row_high, column_low, column_high
                in getattr(sheet, "merged_cells", []) or []
            ]
            if ranges:
                _expand_merges(rows, ranges, first_row=1, first_column=1)
            trimmed, _, _ = _trim_table(rows)
            section = _sheet_pieces(name, index, trimmed, skipped)
            if not section:
                empty.append(name)
                continue
            pieces.extend(section)
            if skipped:
                result.warnings.append(
                    f"лист «{name}»: в текст взяты первые {len(trimmed)} строк, "
                    f"пропущено {skipped}"
                )
    except Exception as error:  # noqa: BLE001 — книга разобрана не полностью
        result.warnings.append(f"книга .xls разобрана не полностью: {_reason(error)}")
    finally:
        try:
            book.release_resources()
        except Exception:  # pragma: no cover — ресурсы уже освобождены
            pass

    result.text = "\n\n".join(pieces)
    if empty:
        result.meta["empty_sheets"] = len(empty)
        result.warnings.append(
            "листов без данных: " + ", ".join(f"«{name}»" for name in empty[:20])
        )
    messages = [line.strip() for line in log.getvalue().splitlines() if line.strip()]
    if messages:
        # xlrd пишет диагностику в журнал, а не в исключения: «нет записи
        # CODEPAGE» означает, что кириллица прочитана как латиница — инженер
        # должен это увидеть, а не гадать, почему в тексте «Ïàðàìåòð».
        result.warnings.append("xlrd сообщает: " + "; ".join(messages[:5]))
    if result.is_empty:
        result.warnings.append("в книге .xls нет данных")
    return result


# ------------------------------------------------------------ CSV и TSV ---

#: Порядок перебора разделителей. Для .csv точка с запятой идёт первой: русский
#: Excel сохраняет CSV именно так, а запятая в таких файлах — десятичная.
_CSV_DELIMITERS = (";", ",", "\t", "|")
_TSV_DELIMITERS = ("\t", ";", ",", "|")


def _sniff_delimiter(text: str, suffix: str) -> str:
    """Разделитель полей: сначала по постоянству числа вхождений, затем Sniffer."""
    candidates = _TSV_DELIMITERS if suffix == ".tsv" else _CSV_DELIMITERS
    lines = [line for line in text.splitlines() if line.strip()][:50]
    for delimiter in candidates:
        counts = [line.count(delimiter) for line in lines]
        if counts and counts[0] > 0 and all(count == counts[0] for count in counts):
            return delimiter
    sample = "\n".join(lines[:20])
    try:
        return csv.Sniffer().sniff(sample, delimiters="".join(candidates)).delimiter
    except csv.Error:
        pass
    return "\t" if suffix == ".tsv" else ","


def _parse_number(text: str) -> float | None:
    cleaned = text.replace(" ", "").replace(" ", "").replace(",", ".")
    cleaned = cleaned.replace("−", "-").lstrip("+")
    if not cleaned or cleaned in ("-", "."):
        return None
    try:
        return float(cleaned)
    except ValueError:
        return None


def _looks_numeric(row: Sequence[str]) -> bool:
    """Похожа ли строка на строку данных, а не на шапку таблицы."""
    filled = [cell.strip() for cell in row if cell.strip()]
    if not filled:
        return False
    return all(_parse_number(cell) is not None for cell in filled)


def convert_csv(path: Path) -> ConvertedDocument:
    """CSV/TSV → Markdown-таблица. Разделитель и кодировка — автоматически."""
    suffix = path.suffix.lower()
    kind = "tsv" if suffix == ".tsv" else "csv"
    result = ConvertedDocument(title=path.stem, meta={"source_format": kind})

    text, encoding, problem = _read_text(path)
    if problem:
        result.warnings.append(problem)
    if encoding:
        result.meta["encoding"] = encoding
    if not text.strip():
        result.warnings.append("файл пуст — таблица не извлечена")
        return result

    delimiter = _sniff_delimiter(text, suffix)
    result.meta["delimiter"] = "\\t" if delimiter == "\t" else delimiter

    rows: List[List[str]] = []
    skipped = 0
    reader = csv.reader(io.StringIO(text, newline=""), delimiter=delimiter)
    try:
        for values in reader:
            if len(rows) >= MAX_SHEET_ROWS:
                if any(value.strip() for value in values):
                    skipped += 1
                continue
            rows.append([_clean_line(value) for value in values])
    except csv.Error as error:
        result.warnings.append(
            f"строка {reader.line_num}: разбор прерван ({_reason(error)}), "
            "дальнейшие строки не взяты"
        )

    trimmed, _, _ = _trim_table(rows)
    if not trimmed:
        result.warnings.append("в файле нет ни одной непустой строки")
        return result

    if _looks_numeric(trimmed[0]):
        # Первая строка — данные, а не шапка: в Markdown-таблице шапка
        # обязательна, поэтому подписываем столбцы по номерам.
        width = max(len(row) for row in trimmed)
        trimmed.insert(0, [f"Столбец {number}" for number in range(1, width + 1)])
        result.meta["header"] = "нет"

    pieces = [page_marker(1)]
    pieces.extend(_table_blocks(trimmed))
    if skipped:
        pieces.append(_limit_note(len(trimmed), skipped))
        result.warnings.append(
            f"в текст взяты первые {len(trimmed)} строк, пропущено {skipped} — "
            "для полной выгрузки разбейте файл на части"
        )
    result.text = "\n\n".join(pieces)
    result.page_count = 1
    return result


# ------------------------------------------------------------ регистрация ---

registry.register(registry.ConverterSpec(
    name="pptx",
    suffixes=(".pptx", ".pptm", ".potx"),
    convert=convert_pptx,
    requires=(
        registry.Requirement(
            "python",
            "pptx",
            pip_hint("python-pptx"),
        ),
    ),
    note="презентации PowerPoint: слайд — раздел, таблицы и заметки докладчика",
))

registry.register(registry.ConverterSpec(
    name="xlsx",
    suffixes=(".xlsx", ".xlsm", ".xltx"),
    convert=convert_xlsx,
    requires=(
        registry.Requirement(
            "python",
            "openpyxl",
            pip_hint("openpyxl"),
        ),
    ),
    note="книги Excel 2007 и новее: лист — раздел, значения формул вычисленные",
))

registry.register(registry.ConverterSpec(
    name="xls",
    suffixes=(".xls",),
    convert=convert_xls,
    requires=(
        registry.Requirement(
            "python",
            "xlrd",
            pip_hint("xlrd") + " (xlrd читает только старый двоичный .xls)",
        ),
    ),
    note="книги Excel 97–2003 (двоичный .xls), выгрузки старых приборов",
))

registry.register(registry.ConverterSpec(
    name="csv",
    suffixes=(".csv", ".tsv"),
    convert=convert_csv,
    note="таблицы CSV и TSV: разделитель и кодировка определяются автоматически",
))
