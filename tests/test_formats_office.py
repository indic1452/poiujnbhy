"""Тесты конвертеров офисных форматов Microsoft: PPTX, XLSX, XLS, CSV/TSV.

Все проверочные файлы собираются программно во временном каталоге: презентации
через python-pptx, книги Excel через openpyxl, CSV — обычной записью байтов. В
репозиторий не кладётся ничего.

Отдельная история с двоичным .xls. Записать его нечем: xlwt в окружении нет, а
LibreOffice здесь поставлен без модуля Calc (только libreoffice-core), поэтому
конвертацией файл тоже не получить. Зато xlrd читает все диалекты BIFF, включая
самый первый, — и BIFF2 достаточно прост, чтобы собрать книгу из записей вручную
(`write_biff2` ниже). Это настоящий файл .xls, который xlrd разбирает штатным
путём, так что конвертер проверяется целиком, а не по частям.
"""

import ast
import io
import os
import re
import shutil
import tempfile
import unittest
import zipfile
from pathlib import Path

import _bootstrap  # noqa: F401

from reportgen import corpus
from reportgen.ingest import convert as convert_module
from reportgen.ingest import registry
from reportgen.ingest.formats import office
from reportgen.ingest.pipeline import chunks_from_markdown

try:  # python-pptx нужен только тестам разбора презентаций
    import pptx
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:  # pragma: no cover — окружение без пакета
    pptx = None

try:
    import openpyxl
except ImportError:  # pragma: no cover — окружение без пакета
    openpyxl = None

try:
    import xlrd
except ImportError:  # pragma: no cover — окружение без пакета
    xlrd = None


#: Однопиксельный PNG: python-pptx требует настоящую картинку, содержимое не важно.
PIXEL_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d4948445200000001000000010806000000"
    "1f15c4890000000a49444154789c6360000002000100ffff0300000600"
    "05570cc7c40000000049454e44ae426082"
)

LONG_TEXT = (
    "Запас на замирания выбран по методике МСЭ-R P.530 с учётом длины пролёта "
    "и профиля местности. Измерения уровня принимаемого сигнала выполнены в "
    "течение суток, отклонения от расчётного значения не превысили допуска, "
    "зафиксированного в протоколе приёмо-сдаточных испытаний линии связи."
)


# ---------------------------------------------------------- вспомогательное ---

def make_presentation(path, slides):
    """Собирает презентацию.

    ``slides`` — список описаний слайда: словарь с ключами ``title``, ``text``
    (список абзацев или пар «абзац, уровень»), ``table`` (список строк),
    ``notes``, ``pictures`` (список подписей; ``None`` — без подписи).
    """
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    title_only = presentation.slide_layouts[5]
    for description in slides:
        layout = title_only if description.get("title") else blank
        slide = presentation.slides.add_slide(layout)
        if description.get("title"):
            slide.shapes.title.text = description["title"]
        top = 2.0
        for paragraph in description.get("text", []):
            text, level = paragraph if isinstance(paragraph, tuple) else (paragraph, 0)
            box = slide.shapes.add_textbox(
                Inches(0.5), Inches(top), Inches(4), Inches(0.8)
            )
            frame = box.text_frame
            frame.text = text
            frame.paragraphs[0].level = level
            top += 0.9
        if description.get("table"):
            rows = description["table"]
            shape = slide.shapes.add_table(
                len(rows), len(rows[0]), Inches(0.5), Inches(top), Inches(6), Inches(1)
            )
            for row_index, row in enumerate(rows):
                for column_index, value in enumerate(row):
                    shape.table.cell(row_index, column_index).text = value
            description["_table"] = shape.table
        for index, caption in enumerate(description.get("pictures", [])):
            picture = slide.shapes.add_picture(
                io.BytesIO(PIXEL_PNG), Inches(6), Inches(1 + index), Inches(0.5), Inches(0.5)
            )
            if caption is not None:
                picture._element._nvXxPr.cNvPr.set("descr", caption)
        if description.get("notes"):
            slide.notes_slide.notes_text_frame.text = description["notes"]
    presentation.save(str(path))
    return presentation


def make_workbook(path, sheets):
    """Собирает книгу Excel. ``sheets`` — список пар «имя листа, словарь ячеек»."""
    workbook = openpyxl.Workbook()
    workbook.remove(workbook.active)
    for name, cells in sheets:
        worksheet = workbook.create_sheet(name)
        for reference, value in cells.items():
            worksheet[reference] = value
    workbook.save(str(path))
    return workbook


def inject_formula(path, cell, formula):
    """Дописывает в ячейку книги формулу, сохраняя уже записанное значение.

    openpyxl умеет записать либо значение, либо формулу, а нужна ячейка с
    формулой И кэшем результата — только такую Excel и сохраняет, и только из
    такой ``data_only=True`` достаёт число. Проще всего дописать элемент ``<f>``
    прямо в XML листа. Возвращает число изменённых ячеек.
    """
    with zipfile.ZipFile(path) as archive:
        items = [(item.filename, archive.read(item.filename)) for item in archive.infolist()]
    pattern = re.compile(r'(<c r="%s"[^>]*>)(<v>)' % re.escape(cell))
    changed = 0
    rewritten = []
    for name, data in items:
        if name.startswith("xl/worksheets/sheet"):
            text, count = pattern.subn(r"\1<f>%s</f>\2" % formula, data.decode("utf-8"))
            changed += count
            data = text.encode("utf-8")
        rewritten.append((name, data))
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, data in rewritten:
            archive.writestr(name, data)
    return changed


def write_biff2(path, rows, sheet_encoding="cp1251"):
    """Записывает книгу .xls в самом старом диалекте BIFF (Excel 2.x).

    Записи BIFF2 устроены как «код (2 байта), длина (2 байта), данные», а из
    типов ячеек нужны только строка и число. Этого хватает, чтобы получить
    настоящий двоичный .xls, который xlrd открывает обычным ``open_workbook``.
    Запись CODEPAGE обязательна: без неё xlrd читает кириллицу как latin-1.
    """
    import struct

    def record(code, data=b""):
        return struct.pack("<HH", code, len(data)) + data

    def label(row, column, text):
        raw = text.encode(sheet_encoding)
        return record(
            0x0004,
            struct.pack("<HH", row, column) + b"\x00\x00\x00" + bytes([len(raw)]) + raw,
        )

    def number(row, column, value):
        return record(
            0x0003,
            struct.pack("<HH", row, column) + b"\x00\x00\x00" + struct.pack("<d", value),
        )

    height = len(rows)
    width = max(len(row) for row in rows) if rows else 0
    stream = record(0x0009, struct.pack("<HH", 0x0000, 0x0010))     # BOF, лист
    stream += record(0x0042, struct.pack("<H", 1251))               # CODEPAGE
    stream += record(0x0000, struct.pack("<HHHH", 0, height, 0, width))  # DIMENSIONS
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            if value is None or value == "":
                continue
            if isinstance(value, (int, float)) and not isinstance(value, bool):
                stream += number(row_index, column_index, float(value))
            else:
                stream += label(row_index, column_index, str(value))
    stream += record(0x000A)                                        # EOF
    Path(path).write_bytes(stream)


def table_blocks(text):
    """Куски текста, являющиеся Markdown-таблицами."""
    return [block for block in text.split("\n\n") if block.startswith("| ")]


class TempCase(unittest.TestCase):
    """Общий временный каталог: проверочные файлы в репозиторий не попадают."""

    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp(prefix="office-"))
        self.addCleanup(shutil.rmtree, self.tmp, ignore_errors=True)


# ----------------------------------------------------------------- PPTX ---

@unittest.skipUnless(pptx is not None, "нужен python-pptx")
class PptxTest(TempCase):

    def convert(self, slides, name="deck.pptx"):
        path = self.tmp / name
        make_presentation(path, slides)
        return office.convert_pptx(path)

    def test_notes_are_extracted(self):
        """Заметки докладчика — часто самое ценное в технической презентации."""
        result = self.convert([
            {"title": "Пролёт", "text": ["Схема"],
             "notes": "Запас 3 дБ согласован с заказчиком"},
        ])
        self.assertIn("Заметки к слайду", result.text)
        self.assertIn("Запас 3 дБ согласован с заказчиком", result.text)

    def test_notes_kept_even_for_slide_without_body(self):
        result = self.convert([
            {"title": "Итоги", "notes": "Повторить измерения после дождя"},
        ])
        self.assertIn("Повторить измерения после дождя", result.text)

    def test_slide_table_is_markdown(self):
        result = self.convert([
            {"title": "Допуски", "table": [
                ["Параметр", "Норма"],
                ["Затухание", "3,5 дБ"],
                ["КБВ", "1,2"],
            ]},
        ])
        blocks = table_blocks(result.text)
        self.assertEqual(len(blocks), 1)
        lines = blocks[0].splitlines()
        self.assertEqual(len(lines), 4)  # шапка, разделитель, две строки данных
        self.assertIn("| Параметр | Норма |", lines[0])
        self.assertIn("| Затухание | 3,5 дБ |", lines[2])

    def test_slide_table_merged_cells_are_expanded(self):
        path = self.tmp / "merged.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1)).table
        table.cell(0, 0).text = "Диапазон частот"
        table.cell(0, 0).merge(table.cell(0, 1))
        table.cell(1, 0).text = "Ku"
        table.cell(1, 1).text = "Ka"
        presentation.save(str(path))
        result = office.convert_pptx(path)
        self.assertIn("| Диапазон частот | Диапазон частот |", result.text)

    def test_page_markers_follow_slide_numbers(self):
        result = self.convert([
            {"title": "Первый", "text": ["раз"]},
            {"title": "Второй", "text": ["два"]},
            {"title": "Третий", "text": ["три"]},
        ])
        self.assertEqual(
            [number for number, _ in convert_module.page_markers(result.text)],
            [1, 2, 3],
        )
        self.assertEqual(result.page_count, 3)

    def test_empty_slides_are_skipped_but_counted(self):
        """Пустой слайд-разделитель в индексе не нужен, но нумерация не сдвигается."""
        result = self.convert([
            {"title": "Первый", "text": ["раз"]},
            {},
            {"title": "Третий", "text": ["три"]},
        ])
        self.assertEqual(result.meta["empty_slides"], 1)
        self.assertNotIn("## Слайд 2", result.text)
        self.assertIn("## Слайд 3", result.text)
        self.assertTrue(any("без содержимого" in item for item in result.warnings))

    def test_picture_alt_text_is_kept(self):
        """Для схемы, вставленной картинкой, alt-текст — единственный текст."""
        result = self.convert([
            {"title": "Тракт", "pictures": ["Схема тракта приёма"]},
        ])
        self.assertIn("![Схема тракта приёма]()", result.text)
        self.assertEqual(result.meta["images"], 1)

    def test_autogenerated_image_name_is_not_used_as_caption(self):
        """python-pptx проставляет описание вида image.png — это не подпись."""
        result = self.convert([{"title": "Тракт", "pictures": [None]}])
        self.assertNotIn("image.png", result.text)
        self.assertIn("![рисунок 1]()", result.text)

    def test_slide_title_becomes_heading(self):
        result = self.convert([{"title": "Параметры пролёта", "text": ["текст"]}])
        self.assertIn("## Слайд 1", result.text)
        self.assertIn("### Параметры пролёта", result.text)
        self.assertEqual(result.title, "Параметры пролёта")

    def test_shapes_are_ordered_top_to_bottom(self):
        """Порядок хранения фигур — это z-order, читать надо сверху вниз."""
        path = self.tmp / "order.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        lower = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
        lower.text_frame.text = "Нижняя рамка"
        upper = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        upper.text_frame.text = "Верхняя рамка"
        presentation.save(str(path))
        result = office.convert_pptx(path)
        self.assertLess(
            result.text.index("Верхняя рамка"), result.text.index("Нижняя рамка")
        )

    def test_bullet_levels_become_list_items(self):
        result = self.convert([
            {"title": "Методика", "text": ["Основание", ("по МСЭ-R P.530", 1)]},
        ])
        self.assertIn("- по МСЭ-R P.530", result.text)

    def test_potx_uses_the_same_converter(self):
        """Шаблон .potx — тот же пакет OOXML, отличается только тип содержимого."""
        source = self.tmp / "deck.pptx"
        make_presentation(source, [{"title": "Шаблон", "text": ["Раздел"]}])
        template = self.tmp / "deck.potx"
        shutil.copy(source, template)
        spec = registry.find(".potx")
        self.assertEqual(spec.name, "pptx")
        result = convert_module.convert_file(template)
        self.assertIn("Раздел", result.text)
        self.assertEqual(result.meta["source_format"], "pptx")

    def test_text_inside_group_is_extracted(self):
        """Схему на слайде часто собирают из сгруппированных подписей."""
        path = self.tmp / "group.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        group = slide.shapes.add_group_shape()
        upper = group.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.6))
        upper.text_frame.text = "Модулятор"
        lower = group.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(0.6))
        lower.text_frame.text = "Усилитель мощности"
        presentation.save(str(path))

        result = office.convert_pptx(path)
        self.assertIn("Модулятор", result.text)
        self.assertIn("Усилитель мощности", result.text)
        self.assertLess(
            result.text.index("Модулятор"), result.text.index("Усилитель мощности")
        )

    def test_chart_title_and_series_are_extracted(self):
        """По подписи диаграммы слайд находится поиском, поэтому она нужна."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        path = self.tmp / "chart.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        data = CategoryChartData()
        data.categories = ["Ku", "Ka"]
        data.add_series("Затухание", (3.5, 5.2))
        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1), Inches(6), Inches(3), data,
        )
        frame.chart.has_title = True
        frame.chart.chart_title.text_frame.text = "Дождевое затухание"
        presentation.save(str(path))

        result = office.convert_pptx(path)
        self.assertIn("**Диаграмма: Дождевое затухание**", result.text)
        self.assertIn("| Категория | Затухание |", result.text)
        self.assertIn("| Ku | 3.5 |", result.text)

    def test_broken_file_gives_warning_not_exception(self):
        path = self.tmp / "broken.pptx"
        path.write_bytes(os.urandom(3072))
        result = office.convert_pptx(path)
        self.assertTrue(result.is_empty)
        self.assertTrue(result.warnings)
        self.assertIn("не удалось открыть презентацию", result.warnings[0])

    def test_presentation_without_slides(self):
        path = self.tmp / "empty.pptx"
        Presentation().save(str(path))
        result = office.convert_pptx(path)
        self.assertTrue(result.is_empty)
        self.assertTrue(any("ни одного слайда" in item for item in result.warnings))


# ----------------------------------------------------------------- XLSX ---

@unittest.skipUnless(openpyxl is not None, "нужен openpyxl")
class XlsxTest(TempCase):

    def test_sheet_becomes_section_with_page_marker(self):
        path = self.tmp / "book.xlsx"
        make_workbook(path, [
            ("Допуски", {"A1": "Параметр", "B1": "Норма", "A2": "КБВ", "B2": 1.2}),
            ("Протокол", {"A1": "Дата", "A2": "05.03.2024"}),
        ])
        result = office.convert_xlsx(path)
        self.assertIn("## Лист «Допуски»", result.text)
        self.assertIn("## Лист «Протокол»", result.text)
        self.assertEqual(
            [number for number, _ in convert_module.page_markers(result.text)], [1, 2]
        )
        self.assertEqual(result.page_count, 2)
        self.assertEqual(result.meta["source_format"], "xlsx")

    def test_cached_formula_value_is_used(self):
        """Нужна цифра допуска, а не строка «=SUM(...)»: читаем кэш значений."""
        path = self.tmp / "formula.xlsx"
        make_workbook(path, [("Расчёт", {"A1": 3, "A2": 4, "A3": 7})])
        self.assertEqual(inject_formula(path, "A3", "SUM(A1:A2)"), 1)

        # Файл действительно содержит формулу, а не просто число.
        raw = openpyxl.load_workbook(str(path), data_only=False)
        self.assertEqual(raw["Расчёт"]["A3"].value, "=SUM(A1:A2)")
        raw.close()

        result = office.convert_xlsx(path)
        self.assertNotIn("SUM", result.text)
        self.assertIn("| 7 |", result.text)

    def test_merged_cells_are_expanded(self):
        path = self.tmp / "merged.xlsx"
        workbook = openpyxl.Workbook()
        worksheet = workbook.active
        worksheet.title = "Шапка"
        worksheet["A1"] = "Диапазон частот"
        worksheet.merge_cells("A1:C1")
        worksheet["A2"] = "Ku"
        worksheet["B2"] = "Ka"
        worksheet["C2"] = "C"
        workbook.save(str(path))
        result = office.convert_xlsx(path)
        self.assertIn(
            "| Диапазон частот | Диапазон частот | Диапазон частот |", result.text
        )

    def test_wide_sheet_is_split_with_repeated_header(self):
        path = self.tmp / "wide.xlsx"
        cells = {}
        for column in range(1, 31):
            letter = openpyxl.utils.get_column_letter(column)
            cells[f"{letter}1"] = f"К{column}"
            cells[f"{letter}2"] = column
        make_workbook(path, [("Широкий", cells)])
        result = office.convert_xlsx(path)

        blocks = table_blocks(result.text)
        self.assertEqual(len(blocks), 2)
        self.assertEqual(len(blocks[0].splitlines()[0].split("|")) - 2, 25)
        self.assertEqual(len(blocks[1].splitlines()[0].split("|")) - 2, 5)
        # Продолжение начинается со своей строки заголовков, а не с голых цифр.
        self.assertTrue(blocks[1].splitlines()[0].startswith("| К26 |"))
        self.assertIn("| --- |", blocks[1].splitlines()[1])
        self.assertIn("### Столбцы 26–30", result.text)

    def test_row_limit_warns_and_reports_how_many_skipped(self):
        path = self.tmp / "huge.xlsx"
        workbook = openpyxl.Workbook()
        worksheet = workbook.active
        worksheet.title = "Выгрузка"
        worksheet.append(["Отсчёт", "Уровень"])
        total = office.MAX_SHEET_ROWS + 100
        for number in range(1, total):
            worksheet.append([number, -70 - number % 5])
        workbook.save(str(path))

        result = office.convert_xlsx(path)
        skipped = total - office.MAX_SHEET_ROWS
        self.assertTrue(any("пропущено" in item for item in result.warnings))
        self.assertTrue(any(str(skipped) in item for item in result.warnings))
        self.assertIn(f"пропущено строк: {skipped}", result.text)
        self.assertEqual(len(table_blocks(result.text)[0].splitlines()),
                         office.MAX_SHEET_ROWS + 1)  # +1 — строка-разделитель

    def test_empty_workbook_warns_instead_of_failing(self):
        path = self.tmp / "empty.xlsx"
        openpyxl.Workbook().save(str(path))
        result = office.convert_xlsx(path)
        self.assertTrue(result.is_empty)
        self.assertTrue(any("листов без данных" in item for item in result.warnings))
        self.assertTrue(any("нет данных" in item for item in result.warnings))

    def test_broken_file_gives_warning_not_exception(self):
        path = self.tmp / "broken.xlsx"
        path.write_bytes(os.urandom(4096))
        result = office.convert_xlsx(path)
        self.assertTrue(result.is_empty)
        self.assertEqual(len(result.warnings), 1)
        self.assertIn("не удалось открыть книгу Excel", result.warnings[0])

    def test_empty_edges_are_trimmed(self):
        """Таблица, начатая с D5, не должна тащить за собой пустые поля."""
        path = self.tmp / "offset.xlsx"
        make_workbook(path, [
            ("Смещённый", {"D5": "Параметр", "E5": "Норма", "D6": "КБВ", "E6": 1.2}),
        ])
        result = office.convert_xlsx(path)
        block = table_blocks(result.text)[0]
        self.assertTrue(block.splitlines()[0].startswith("| Параметр | Норма |"))
        self.assertEqual(len(block.splitlines()), 3)

    def test_numbers_and_dates_are_readable(self):
        import datetime

        path = self.tmp / "types.xlsx"
        make_workbook(path, [("Значения", {
            "A1": "Малое", "B1": "Большое", "C1": "Дата", "D1": "Время", "E1": "Флаг",
            "A2": 0.000001234,
            "B2": 12345678901234.0,
            "C2": datetime.datetime(2024, 3, 5),
            "D2": datetime.datetime(2024, 3, 5, 12, 30),
            "E2": True,
        })])
        result = office.convert_xlsx(path)
        self.assertIn("0.000001234", result.text)
        self.assertNotIn("e-06", result.text)
        self.assertNotIn("E+", result.text)
        self.assertIn("12345678901234", result.text)
        self.assertIn("2024-03-05", result.text)
        self.assertIn("2024-03-05 12:30:00", result.text)
        self.assertIn("истина", result.text)

    def test_macro_and_template_books_use_the_same_converter(self):
        source = self.tmp / "book.xlsx"
        make_workbook(source, [("Данные", {"A1": "Параметр", "A2": "КБВ"})])
        for suffix in (".xlsm", ".xltx"):
            with self.subTest(suffix=suffix):
                copy = self.tmp / f"book{suffix}"
                shutil.copy(source, copy)
                self.assertEqual(registry.find(suffix).name, "xlsx")
                result = convert_module.convert_file(copy)
                self.assertIn("КБВ", result.text)

    def test_table_is_not_broken_by_the_chunker(self):
        """Таблица допусков и есть смысл документа — рвать её пополам нельзя."""
        path = self.tmp / "tolerance.xlsx"
        cells = {"A1": "Параметр", "B1": "Норма"}
        for row in range(2, 40):
            cells[f"A{row}"] = f"Параметр канала номер {row}"
            cells[f"B{row}"] = row * 1.5
        make_workbook(path, [("Допуски", cells)])
        result = office.convert_xlsx(path)
        pieces = [piece for _, piece in corpus.split_document(result.text)]
        whole = [
            piece for piece in pieces
            if "| Параметр | Норма |" in piece and "Параметр канала номер 39" in piece
        ]
        self.assertEqual(len(whole), 1)

    def test_hidden_sheets_are_read_and_reported(self):
        """Скрытый лист часто и есть источник данных — берём, но предупреждаем."""
        path = self.tmp / "hidden.xlsx"
        workbook = openpyxl.Workbook()
        visible = workbook.active
        visible.title = "Отчёт"
        visible["A1"] = "Сводка"
        source = workbook.create_sheet("Исходные")
        source["A1"] = "Отсчёт"
        source["A2"] = 42
        source.sheet_state = "hidden"
        workbook.save(str(path))

        result = office.convert_xlsx(path)
        self.assertIn("## Лист «Исходные»", result.text)
        self.assertTrue(any("скрытые листы" in item for item in result.warnings))

    def test_password_protected_book_gets_a_useful_hint(self):
        """Зашифрованная книга снаружи — контейнер OLE, а не zip."""
        path = self.tmp / "protected.xlsx"
        path.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + os.urandom(1024))
        result = office.convert_xlsx(path)
        self.assertTrue(result.is_empty)
        self.assertIn("паролем", result.warnings[0])

    def test_pipe_in_cell_does_not_break_the_table(self):
        path = self.tmp / "pipe.xlsx"
        make_workbook(path, [("Данные", {"A1": "Режим", "A2": "приём|передача"})])
        result = office.convert_xlsx(path)
        self.assertIn("приём\\|передача", result.text)


# ------------------------------------------------------------------ XLS ---

@unittest.skipUnless(xlrd is not None, "нужен xlrd")
class XlsTest(TempCase):

    def test_binary_workbook_is_read(self):
        path = self.tmp / "старая.xls"
        write_biff2(path, [
            ["Параметр", "Значение"],
            ["Затухание", 3.5],
            ["КБВ", 1.2],
        ])
        result = office.convert_xls(path)
        self.assertEqual(result.meta["source_format"], "xls")
        self.assertEqual(result.page_count, 1)
        self.assertIn("| Параметр | Значение |", result.text)
        self.assertIn("| Затухание | 3.5 |", result.text)
        self.assertEqual(
            [number for number, _ in convert_module.page_markers(result.text)], [1]
        )

    def test_dispatch_through_registry(self):
        path = self.tmp / "прибор.xls"
        write_biff2(path, [["Канал", "Уровень"], ["1", -70.0]])
        spec = registry.find(".xls")
        self.assertEqual(spec.name, "xls")
        self.assertEqual([item.name for item in spec.requires], ["xlrd"])
        result = convert_module.convert_file(path)
        self.assertIn("Канал", result.text)

    def test_broken_file_gives_warning_not_exception(self):
        path = self.tmp / "broken.xls"
        path.write_bytes(os.urandom(2048))
        result = office.convert_xls(path)
        self.assertTrue(result.is_empty)
        self.assertTrue(any("не удалось открыть" in item for item in result.warnings))


# ------------------------------------------------------------ CSV и TSV ---

class CsvTest(TempCase):

    def test_semicolon_and_cp1251(self):
        """Русский Excel сохраняет CSV с точкой с запятой и в cp1251."""
        path = self.tmp / "измерения.csv"
        path.write_bytes(
            "Параметр;Значение;Дата\r\n"
            "Затухание;3,5;2024-03-05\r\n"
            "КБВ;1,2;2024-03-06\r\n".encode("cp1251")
        )
        result = office.convert_csv(path)
        self.assertEqual(result.meta["encoding"], "cp1251")
        self.assertEqual(result.meta["delimiter"], ";")
        self.assertIn("| Параметр | Значение | Дата |", result.text)
        self.assertIn("| Затухание | 3,5 | 2024-03-05 |", result.text)
        self.assertEqual(len(table_blocks(result.text)[0].splitlines()), 4)

    def test_tsv(self):
        path = self.tmp / "каналы.tsv"
        path.write_text("Канал\tУровень\n1\t-70\n2\t-72\n", encoding="utf-8")
        result = office.convert_csv(path)
        self.assertEqual(result.meta["source_format"], "tsv")
        self.assertEqual(result.meta["delimiter"], "\\t")
        self.assertIn("| Канал | Уровень |", result.text)
        self.assertIn("| 2 | -72 |", result.text)

    def test_single_column(self):
        """csv.Sniffer на одном столбце разделитель не находит — это не ошибка."""
        path = self.tmp / "частоты.csv"
        path.write_text("Частота\n1200\n1300\n", encoding="utf-8")
        result = office.convert_csv(path)
        self.assertIn("| Частота |", result.text)
        self.assertIn("| 1200 |", result.text)
        self.assertEqual(len(table_blocks(result.text)[0].splitlines()), 4)
        self.assertFalse(result.warnings)

    def test_numeric_first_row_is_data_not_header(self):
        path = self.tmp / "без-шапки.csv"
        path.write_text("1,2,3\n4,5,6\n", encoding="utf-8")
        result = office.convert_csv(path)
        self.assertEqual(result.meta["header"], "нет")
        self.assertIn("| Столбец 1 | Столбец 2 | Столбец 3 |", result.text)
        self.assertIn("| 1 | 2 | 3 |", result.text)

    def test_page_marker_is_present(self):
        path = self.tmp / "выгрузка.csv"
        path.write_text("Канал,Уровень\n1,-70\n", encoding="utf-8")
        result = office.convert_csv(path)
        self.assertEqual(
            [number for number, _ in convert_module.page_markers(result.text)], [1]
        )
        self.assertEqual(result.page_count, 1)

    def test_empty_file_warns(self):
        path = self.tmp / "пусто.csv"
        path.write_text("", encoding="utf-8")
        result = office.convert_csv(path)
        self.assertTrue(result.is_empty)
        self.assertTrue(any("пуст" in item for item in result.warnings))

    def test_row_limit_warns(self):
        path = self.tmp / "длинная.csv"
        total = office.MAX_SHEET_ROWS + 50
        lines = ["Отсчёт,Уровень"]
        lines += [f"{number},-7{number % 10}" for number in range(total)]
        path.write_text("\n".join(lines) + "\n", encoding="utf-8")
        result = office.convert_csv(path)
        skipped = total + 1 - office.MAX_SHEET_ROWS
        self.assertTrue(any(f"пропущено {skipped}" in item for item in result.warnings))
        self.assertIn(f"пропущено строк: {skipped}", result.text)

    def test_quoted_field_with_newline(self):
        path = self.tmp / "цитата.csv"
        path.write_text(
            'Параметр,Примечание\nКБВ,"первая строка\nвторая строка"\n', encoding="utf-8"
        )
        result = office.convert_csv(path)
        self.assertIn("первая строка <br> вторая строка", result.text)
        self.assertEqual(len(table_blocks(result.text)[0].splitlines()), 3)

    def test_dispatch_through_registry(self):
        path = self.tmp / "данные.tsv"
        path.write_text("Канал\tУровень\n1\t-70\n", encoding="utf-8")
        result = convert_module.convert_file(path)
        self.assertIn("| Канал | Уровень |", result.text)


# ------------------------------------------------------------- реестр ----

class RegistryTest(unittest.TestCase):

    def test_all_office_formats_are_registered(self):
        expected = {
            ".pptx": "pptx", ".pptm": "pptx", ".potx": "pptx",
            ".xlsx": "xlsx", ".xlsm": "xlsx", ".xltx": "xlsx",
            ".xls": "xls",
            ".csv": "csv", ".tsv": "csv",
        }
        for suffix, name in expected.items():
            with self.subTest(suffix=suffix):
                spec = registry.find(suffix)
                self.assertIsNotNone(spec, f"формат {suffix} не зарегистрирован")
                self.assertEqual(spec.name, name)

    def test_requirements_carry_install_hints(self):
        """На машине заказчика пакетов может не быть — подсказка обязана работать.

        Прежняя подсказка звучала «pip install openpyxl; в Windows — py -m pip
        install openpyxl» и была негодной дважды: интернета в контуре нет, а
        py — системный запускатель, и пакет уехал бы мимо окружения, в котором
        работает приложение. Теперь называется тот самый интерпретатор и
        установка без сети.
        """
        import sys

        for name, package in (("pptx", "pptx"), ("xlsx", "openpyxl"), ("xls", "xlrd")):
            with self.subTest(name=name):
                spec = next(item for item in registry.all_specs() if item.name == name)
                requirement = next(
                    item for item in spec.requires if item.name == package
                )
                self.assertEqual(requirement.kind, "python")
                self.assertIn("pip install", requirement.hint)
                self.assertIn("--no-index", requirement.hint,
                              "подсказка уводит в интернет, которого нет")
                self.assertIn(sys.executable, requirement.hint,
                              "подсказка не называет интерпретатор приложения")
                self.assertNotIn("py -m pip", requirement.hint,
                                 "py — другой интерпретатор: пакет уедет мимо окружения")
                self.assertTrue(spec.note)

    def test_csv_needs_nothing_extra(self):
        spec = next(item for item in registry.all_specs() if item.name == "csv")
        self.assertEqual(spec.requires, ())
        self.assertTrue(spec.is_available())

    def test_module_import_needs_no_third_party_packages(self):
        """Один недостающий пакет не должен лишать систему остальных форматов."""
        source = Path(office.__file__).read_text(encoding="utf-8")
        tree = ast.parse(source)
        allowed = set(
            "ast csv datetime decimal io os posixpath re shutil typing "
            "xml zipfile pathlib __future__".split()
        )
        for node in tree.body:
            if isinstance(node, ast.Import):
                for alias in node.names:
                    self.assertIn(alias.name.split(".")[0], allowed, alias.name)
            elif isinstance(node, ast.ImportFrom):
                if node.level:  # относительный импорт внутри пакета
                    continue
                self.assertIn(node.module.split(".")[0], allowed, node.module)


# --------------------------------------------------------- интеграция ----

@unittest.skipUnless(pptx is not None and openpyxl is not None, "нужны pptx и openpyxl")
class ChunkingTest(TempCase):

    def test_slide_number_reaches_chunk_meta(self):
        path = self.tmp / "deck.pptx"
        make_presentation(path, [
            {"title": "Первый пролёт", "text": [LONG_TEXT]},
            {"title": "Второй пролёт", "text": [LONG_TEXT]},
        ])
        result = office.convert_pptx(path)
        chunks = chunks_from_markdown(result.text, "presentations/deck", "literature")
        pages = sorted({chunk.meta.get("page") for chunk in chunks})
        self.assertEqual(pages, [1, 2])

    def test_sheet_number_reaches_chunk_meta(self):
        path = self.tmp / "book.xlsx"
        make_workbook(path, [
            ("Первый", {f"A{row}": f"{LONG_TEXT} {row}" for row in range(1, 4)}),
            ("Второй", {f"A{row}": f"{LONG_TEXT} {row}" for row in range(1, 4)}),
        ])
        result = office.convert_xlsx(path)
        chunks = chunks_from_markdown(result.text, "standards/book", "standards")
        pages = sorted({chunk.meta.get("page") for chunk in chunks})
        self.assertEqual(pages, [1, 2])
        self.assertTrue(
            any("Лист «Второй»" in chunk.title_path for chunk in chunks)
        )


if __name__ == "__main__":  # pragma: no cover — ручной запуск
    unittest.main()
