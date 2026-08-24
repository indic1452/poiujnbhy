"""Тесты старых форматов: путь через headless LibreOffice и разбор RTF.

Как делаются проверочные файлы:

* **RTF** — текстовый формат, он пишется строкой прямо в тесте: так видно, что
  именно проверяется (стиль заголовка, ячейки таблицы, разрыв страницы), и
  ничего стороннего для этого не нужно.
* **DOC и PPT** — двоичные форматы девяностых, записать их нечем. По замыслу
  они собираются так: python-docx и python-pptx делают современный файл, а сам
  LibreOffice сохраняет его в старом формате (``--convert-to doc``). Такие
  тесты обёрнуты в ``skipUnless``: в окружении сборки стоит только
  ``libreoffice-core``, без модулей Writer/Calc/Impress, поэтому фильтров
  документов нет и настоящая конвертация недоступна (проба
  :func:`_libreoffice_converts` это выясняет один раз при импорте).
* **Подставной soffice** — чтобы вся обвязка (ключи запуска, отдельный профиль,
  таймаут, разбор результата, удаление временных файлов, параллельные запуски)
  проверялась и без установленного LibreOffice, рядом кладётся исполняемый
  скрипт с тем же интерфейсом командной строки. Он кладёт в целевой каталог
  заранее подготовленный настоящий DOCX/PPTX/XLSX/PDF — то есть проверяется всё
  ровно до вызова LibreOffice и всё после него. Скрипт исполняемый, поэтому
  такие тесты идут только на POSIX.

Ничего в репозиторий не пишется: все файлы живут во временных каталогах.
"""

import json
import os
import shutil
import sys
import tempfile
import threading
import unittest
from pathlib import Path
from unittest import mock

import _bootstrap  # noqa: F401

from reportgen.ingest import registry
from reportgen.ingest.convert import convert_file
from reportgen.ingest.formats import legacy
from reportgen.ingest.pipeline import chunks_from_markdown

try:  # python-docx нужен для сборки промежуточного DOCX
    import docx
except ImportError:  # pragma: no cover — окружение без пакета
    docx = None

try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:  # pragma: no cover — окружение без пакета
    pptx = None

try:
    import openpyxl
except ImportError:  # pragma: no cover — окружение без пакета
    openpyxl = None

try:
    import pymupdf
except ImportError:  # pragma: no cover — окружение без пакета
    pymupdf = None

try:
    import striprtf
except ImportError:  # pragma: no cover — окружение без пакета
    striprtf = None

try:
    import xlrd
except ImportError:  # pragma: no cover — окружение без пакета
    xlrd = None


POSIX = os.name == "posix"

#: Абзац нарочно длиннее corpus.MIN_CHARS: короткий фрагмент слой приёма
#: присоединяет к предыдущему чанку, и номер страницы до метаданных не доезжает.
LONG_TEXT = (
    "Запас на замирания выбран по методике МСЭ-R P.530 с учётом длины пролёта, "
    "профиля местности и требуемого коэффициента готовности линии связи. "
    "Уровень принимаемого сигнала измерялся в течение суток, отклонения от "
    "расчётного значения не превысили допуска приёмо-сдаточных испытаний."
)


# ---------------------------------------------------------- вспомогательное ---

def make_docx(directory, name="протокол.docx"):
    """Современный DOCX с заголовком, абзацем и таблицей допусков."""
    document = docx.Document()
    document.add_heading("Протокол измерений", level=1)
    document.add_paragraph(LONG_TEXT)
    document.add_heading("Допуски", level=2)
    table = document.add_table(rows=3, cols=2)
    values = (
        ("Параметр", "Значение"),
        ("EVM, %", "3,2"),
        ("Уровень приёма, дБм", "минус 42"),
    )
    for row, pair in zip(table.rows, values, strict=True):
        row.cells[0].text = pair[0]
        row.cells[1].text = pair[1]
    path = Path(directory) / name
    document.save(str(path))
    return path


def make_pptx(directory, name="доклад.pptx"):
    """Презентация из двух слайдов с заголовками, текстом и заметками."""
    presentation = Presentation()
    layout = presentation.slide_layouts[5]
    # Текст слайда длинный не для красоты: короткий чанк слой приёма
    # присоединяет к предыдущему, и номер слайда до метаданных не доезжает.
    for number, (title, body) in enumerate(
        (
            ("Состояние линии", f"Готовность 99,95 процента. {LONG_TEXT}"),
            ("План работ", f"Замена ODU на пролёте 4. {LONG_TEXT}"),
        ),
        start=1,
    ):
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6), Inches(1))
        box.text_frame.text = body
        slide.notes_slide.notes_text_frame.text = f"Заметка к слайду {number}"
    path = Path(directory) / name
    presentation.save(str(path))
    return path


def make_xlsx(directory, name="бюджет.xlsx"):
    """Книга Excel с одним листом и таблицей."""
    book = openpyxl.Workbook()
    sheet = book.active
    sheet.title = "Бюджет канала"
    sheet.append(["Параметр", "Значение"])
    sheet.append(["ЭИИМ, дБВт", 52.4])
    path = Path(directory) / name
    book.save(str(path))
    return path


#: Встроенные шрифты PDF (Base-14) кириллицы не содержат, поэтому для
#: проверочного файла нужен настоящий TTF. Если ни одного нет — тест на PUB
#: пропускается, а не проверяет расставленные точки вместо букв.
CYRILLIC_FONTS = (
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "/usr/share/fonts/TTF/DejaVuSans.ttf",
    "C:\\Windows\\Fonts\\arial.ttf",
)


def cyrillic_font():
    for candidate in CYRILLIC_FONTS:
        if Path(candidate).is_file():
            return candidate
    return None


def make_pdf(directory, name="буклет.pdf"):
    """PDF из двух страниц — им притворяется результат конвертации Publisher."""
    font = cyrillic_font()
    document = pymupdf.open()
    for number, text in enumerate(("Земная станция связи", "Состав оборудования"), start=1):
        page = document.new_page()
        page.insert_text((72, 100), f"{text} {number}", fontsize=18,
                         fontname="tt", fontfile=font)
        # Именно textbox: insert_text не переносит строки, текст обрезался бы
        # по краю листа и страница вышла бы короче порога нарезки на чанки.
        page.insert_textbox(pymupdf.Rect(72, 120, 520, 400),
                            f"{LONG_TEXT} Страница {number}.", fontsize=11,
                            fontname="tt", fontfile=font)
    path = Path(directory) / name
    document.save(str(path))
    document.close()
    return path


def make_fake_soffice(directory, *, payload=None, sleep=0.0, log=None):
    """Подставной ``soffice`` с интерфейсом настоящего.

    Кладёт ``payload`` в каталог из ``--outdir`` под именем исходного файла с
    расширением из ``--convert-to``; без ``payload`` ведёт себя как LibreOffice
    на нечитаемом файле — печатает ошибку и выходит с кодом 0.
    """
    config = {
        "payload": str(payload) if payload else "",
        "sleep": float(sleep),
        "log": str(log) if log else "",
    }
    script = Path(directory) / "soffice"
    script.write_text(
        "#!" + sys.executable + "\n"
        "import json, os, shutil, sys, time\n"
        "CONFIG = json.loads(r'''" + json.dumps(config) + "''')\n"
        "argv = sys.argv[1:]\n"
        "if CONFIG['log']:\n"
        "    with open(CONFIG['log'], 'a', encoding='utf-8') as handle:\n"
        "        handle.write(json.dumps(argv, ensure_ascii=False) + '\\n')\n"
        "if CONFIG['sleep']:\n"
        "    time.sleep(CONFIG['sleep'])\n"
        "target = argv[argv.index('--convert-to') + 1].split(':')[0]\n"
        "outdir = argv[argv.index('--outdir') + 1]\n"
        "source = argv[-1]\n"
        "if CONFIG['payload']:\n"
        "    name = os.path.splitext(os.path.basename(source))[0] + '.' + target\n"
        "    shutil.copyfile(CONFIG['payload'], os.path.join(outdir, name))\n"
        "else:\n"
        "    sys.stdout.write('Error: source file could not be loaded\\n')\n"
        "sys.exit(0)\n",
        encoding="utf-8",
    )
    script.chmod(0o755)
    return script


def use_fake_soffice(script):
    """Подменяет поиск LibreOffice на подставной скрипт."""
    return mock.patch.object(legacy, "soffice_binary", lambda **kwargs: str(script))


def no_soffice():
    """Окружение, в котором LibreOffice не установлен."""
    return mock.patch.object(legacy, "soffice_binary", lambda **kwargs: None)


def soffice_leftovers():
    """Временные каталоги конвертации, оставшиеся в системном temp."""
    return sorted(Path(tempfile.gettempdir()).glob("reportgen-soffice-*"))


def read_log(path):
    lines = Path(path).read_text(encoding="utf-8").splitlines()
    return [json.loads(line) for line in lines if line.strip()]


def option_value(argv, key):
    return argv[argv.index(key) + 1]


def _libreoffice_converts(builder, target):
    """Умеет ли здешний LibreOffice открыть формат и сохранить в ``target``.

    Наличие ``soffice`` в PATH ничего не гарантирует: пакет
    ``libreoffice-core`` ставится без фильтров Writer, Calc и Impress, и тогда
    любая конвертация документа кончается «source file could not be loaded».
    Проба делается один раз при импорте модуля — она занимает секунду.
    """
    if builder is None or legacy.soffice_binary() is None:
        return False
    workdir = Path(tempfile.mkdtemp(prefix="reportgen-probe-"))
    try:
        source = builder(workdir)
        with legacy.convert_with_soffice(source, target, timeout=60) as produced:
            return produced.path is not None
    except Exception:  # pragma: no cover — окружение без LibreOffice
        return False
    finally:
        shutil.rmtree(str(workdir), ignore_errors=True)


LIBREOFFICE_WRITER = _libreoffice_converts(make_docx if docx else None, "doc")
LIBREOFFICE_IMPRESS = _libreoffice_converts(make_pptx if pptx else None, "ppt")


# --------------------------------------------------------- поиск LibreOffice ---

class SofficeBinaryTest(unittest.TestCase):
    """Поиск исполняемого файла: PATH, переменная окружения, каталоги Windows."""

    def setUp(self):
        # На рабочей машине путь к LibreOffice может быть задан переменной
        # окружения — она сильнее любого поиска, и тесты поиска от неё чистим.
        environment = mock.patch.dict(os.environ, {legacy.SOFFICE_ENV_VAR: ""})
        environment.start()
        self.addCleanup(environment.stop)
        legacy.soffice_binary(refresh=True)

    def tearDown(self):
        legacy.soffice_binary(refresh=True)

    def test_находит_soffice_в_path(self):
        with mock.patch.object(legacy, "_windows", lambda: False), \
                mock.patch.object(legacy.shutil, "which",
                                  lambda name: "/usr/bin/soffice" if name == "soffice" else None):
            self.assertEqual(legacy.soffice_binary(refresh=True), "/usr/bin/soffice")

    def test_кэширует_результат_поиска(self):
        calls = []

        def which(name):
            calls.append(name)
            return "/usr/bin/soffice" if name == "soffice" else None

        with mock.patch.object(legacy, "_windows", lambda: False), \
                mock.patch.object(legacy.shutil, "which", which):
            first = legacy.soffice_binary(refresh=True)
            second = legacy.soffice_binary()
            third = legacy.soffice_binary()
        self.assertEqual(first, second)
        self.assertEqual(second, third)
        self.assertEqual(calls, ["soffice"], "поиск должен выполняться один раз")

    def test_переменная_окружения_важнее_path(self):
        directory = Path(tempfile.mkdtemp())
        self.addCleanup(shutil.rmtree, str(directory), True)
        portable = directory / "soffice"
        portable.write_text("", encoding="utf-8")
        portable.chmod(0o755)
        with mock.patch.dict(os.environ, {legacy.SOFFICE_ENV_VAR: str(portable)}), \
                mock.patch.object(legacy, "_windows", lambda: False), \
                mock.patch.object(legacy.shutil, "which", lambda name: "/usr/bin/soffice"):
            self.assertEqual(legacy.soffice_binary(refresh=True), str(portable))

    @unittest.skipUnless(POSIX, "проверка прав на исполнение осмысленна только на POSIX")
    def test_переменная_окружения_принимает_каталог_установки(self):
        directory = Path(tempfile.mkdtemp())
        self.addCleanup(shutil.rmtree, str(directory), True)
        program = directory / "program"
        program.mkdir()
        binary = program / "soffice"
        binary.write_text("", encoding="utf-8")
        binary.chmod(0o755)
        with mock.patch.dict(os.environ, {legacy.SOFFICE_ENV_VAR: str(directory)}), \
                mock.patch.object(legacy, "_windows", lambda: False), \
                mock.patch.object(legacy.shutil, "which", lambda name: None):
            self.assertEqual(legacy.soffice_binary(refresh=True), str(binary))

    def test_стандартные_каталоги_windows(self):
        """В Windows LibreOffice не прописан в PATH — ищем в Program Files."""
        directory = Path(tempfile.mkdtemp())
        self.addCleanup(shutil.rmtree, str(directory), True)
        program = directory / "LibreOffice" / "program"
        program.mkdir(parents=True)
        console = program / "soffice.com"
        console.write_text("", encoding="utf-8")
        (program / "soffice.exe").write_text("", encoding="utf-8")
        environment = {"ProgramFiles": str(directory)}
        with mock.patch.dict(os.environ, environment, clear=False), \
                mock.patch.object(legacy, "_windows", lambda: True), \
                mock.patch.object(legacy.shutil, "which", lambda name: None):
            found = legacy.soffice_binary(refresh=True)
        # Консольный soffice.com дожидается конца конвертации, soffice.exe — нет.
        self.assertEqual(found, str(console))

    def test_без_libreoffice_возвращает_none(self):
        with mock.patch.object(legacy, "_windows", lambda: False), \
                mock.patch.object(legacy.shutil, "which", lambda name: None), \
                mock.patch.object(legacy, "_candidate_paths", list), \
                mock.patch.dict(os.environ, {legacy.SOFFICE_ENV_VAR: ""}):
            self.assertIsNone(legacy.soffice_binary(refresh=True))


# ------------------------------------------------------------ запуск soffice ---

@unittest.skipUnless(POSIX, "подставной soffice — исполняемый скрипт, только POSIX")
class ConvertWithSofficeTest(unittest.TestCase):
    """Обвязка запуска: ключи, профиль, разбор результата, уборка, таймаут."""

    def setUp(self):
        self.workdir = Path(tempfile.mkdtemp(prefix="reportgen-test-"))
        self.addCleanup(shutil.rmtree, str(self.workdir), True)
        self.source = self.workdir / "исходный.doc"
        self.source.write_bytes(b"old binary document")
        self.payload = self.workdir / "payload.docx"
        self.payload.write_bytes(b"PK\x03\x04 payload")
        self.log = self.workdir / "argv.log"

    def test_отдаёт_преобразованный_файл(self):
        script = make_fake_soffice(self.workdir, payload=self.payload, log=self.log)
        with use_fake_soffice(script):
            with legacy.convert_with_soffice(self.source, "docx") as produced:
                self.assertIsNotNone(produced.path)
                self.assertEqual(produced.path.suffix, ".docx")
                self.assertEqual(produced.path.read_bytes(), self.payload.read_bytes())
                self.assertEqual(produced.via, "libreoffice→docx")
                inside = produced.path
        self.assertFalse(inside.exists(), "временный файл должен удаляться после блока")

    def test_ключи_запуска_и_отдельный_профиль(self):
        script = make_fake_soffice(self.workdir, payload=self.payload, log=self.log)
        with use_fake_soffice(script):
            with legacy.convert_with_soffice(self.source, "docx") as produced:
                self.assertIsNotNone(produced.path)
        argv = read_log(self.log)[0]
        for key in ("--headless", "--norestore", "--nolockcheck", "--nodefault",
                    "--nofirststartwizard"):
            self.assertIn(key, argv)
        self.assertEqual(option_value(argv, "--convert-to"), "docx")
        profile = [item for item in argv if item.startswith("-env:UserInstallation=")]
        self.assertEqual(len(profile), 1, "профиль обязателен: общий ломает параллельный запуск")
        self.assertTrue(profile[0].startswith("-env:UserInstallation=file:///"))
        # Исходник копируется во временный каталог: библиотека заказчика не
        # должна обрастать файлами блокировки LibreOffice.
        self.assertNotEqual(Path(argv[-1]), self.source)
        # Имя копии — латиницей: LibreOffice в части сборок не открывает пути
        # с кириллицей, а в библиотеке заказчика так названо почти всё.
        copied = Path(argv[-1])
        self.assertTrue(str(copied).isascii(), f"путь для LibreOffice не латинский: {copied}")
        self.assertEqual(copied.suffix, self.source.suffix)

    def test_временные_файлы_удаляются(self):
        script = make_fake_soffice(self.workdir, payload=self.payload)
        before = soffice_leftovers()
        with use_fake_soffice(script):
            with legacy.convert_with_soffice(self.source, "docx") as produced:
                self.assertIsNotNone(produced.path)
        self.assertEqual(soffice_leftovers(), before)

    def test_таймаут_из_константы(self):
        """LibreOffice умеет висеть вечно: ждём не дольше SOFFICE_TIMEOUT."""
        script = make_fake_soffice(self.workdir, payload=self.payload, sleep=30)
        before = soffice_leftovers()
        with use_fake_soffice(script), mock.patch.object(legacy, "SOFFICE_TIMEOUT", 0.5):
            with legacy.convert_with_soffice(self.source, "docx") as produced:
                self.assertIsNone(produced.path)
                self.assertIn("0.5", produced.warning)
                self.assertIn("LibreOffice", produced.warning)
        self.assertEqual(soffice_leftovers(), before, "по таймауту тоже надо убирать за собой")

    def test_soffice_ничего_не_создал(self):
        """Код возврата 0 без файла — обычное поведение на нечитаемом документе."""
        script = make_fake_soffice(self.workdir, payload=None)
        with use_fake_soffice(script):
            with legacy.convert_with_soffice(self.source, "docx") as produced:
                self.assertIsNone(produced.path)
                self.assertIn("DOCX", produced.warning)
                self.assertIn("could not be loaded", produced.warning)

    def test_несуществующий_исполняемый_файл(self):
        missing = self.workdir / "нет-такого-soffice"
        with mock.patch.object(legacy, "soffice_binary", lambda **kwargs: str(missing)):
            with legacy.convert_with_soffice(self.source, "docx") as produced:
                self.assertIsNone(produced.path)
                self.assertIn("не удалось запустить LibreOffice", produced.warning)


class ConvertWithoutSofficeTest(unittest.TestCase):
    """Без LibreOffice система объясняет, что доложить в комплект."""

    def setUp(self):
        self.workdir = Path(tempfile.mkdtemp(prefix="reportgen-test-"))
        self.addCleanup(shutil.rmtree, str(self.workdir), True)

    def test_конвертация_без_libreoffice_не_падает(self):
        source = self.workdir / "методика.doc"
        source.write_bytes(b"\xd0\xcf\x11\xe0 old word document")
        with no_soffice():
            with legacy.convert_with_soffice(source, "docx") as produced:
                self.assertIsNone(produced.path)
                self.assertIn("LibreOffice", produced.warning)
            result = legacy.convert_doc(source)
        self.assertEqual(result.text, "")
        self.assertTrue(result.warnings)
        warning = " ".join(result.warnings)
        self.assertIn("LibreOffice", warning)
        self.assertIn("Program Files", warning, "подсказка нужна и для Windows")
        self.assertEqual(result.meta["source_format"], "doc")


# ------------------------------------------------------------------ DOC и PPT ---

@unittest.skipUnless(POSIX and docx is not None, "нужен python-docx и POSIX")
class LegacyDocTest(unittest.TestCase):
    """DOC → DOCX → штатный разбор DOCX."""

    def setUp(self):
        self.workdir = Path(tempfile.mkdtemp(prefix="reportgen-test-"))
        self.addCleanup(shutil.rmtree, str(self.workdir), True)

    def test_текст_и_таблица_из_doc(self):
        payload = make_docx(self.workdir)
        script = make_fake_soffice(self.workdir, payload=payload, log=self.workdir / "argv.log")
        source = self.workdir / "протокол.doc"
        source.write_bytes(b"\xd0\xcf\x11\xe0 old word document")
        with use_fake_soffice(script):
            result = legacy.convert_doc(source)
        self.assertEqual(result.warnings, [])
        self.assertIn("# Протокол измерений", result.text)
        self.assertIn("## Допуски", result.text)
        self.assertIn("| Параметр | Значение |", result.text)
        self.assertIn("| EVM, % | 3,2 |", result.text)
        self.assertIn("| Уровень приёма, дБм | минус 42 |", result.text)
        self.assertEqual(result.meta["source_format"], "doc")
        self.assertEqual(result.meta["converted_via"], "libreoffice→docx")
        argv = read_log(self.workdir / "argv.log")[0]
        self.assertEqual(option_value(argv, "--convert-to"), "docx")

    @unittest.skipUnless(LIBREOFFICE_WRITER,
                         "здешний LibreOffice без модуля Writer: .doc не собрать и не прочитать")
    def test_настоящий_doc_через_libreoffice(self):
        """Полный круг: DOCX → (LibreOffice) → DOC → конвертер → Markdown."""
        modern = make_docx(self.workdir)
        legacy_path = self.workdir / "протокол.doc"
        with legacy.convert_with_soffice(modern, "doc") as produced:
            self.assertIsNotNone(produced.path, produced.warning)
            shutil.copyfile(str(produced.path), str(legacy_path))
        result = legacy.convert_doc(legacy_path)
        self.assertIn("Протокол измерений", result.text)
        self.assertIn("| Параметр | Значение |", result.text)
        self.assertIn("3,2", result.text)
        self.assertEqual(result.meta["converted_via"], "libreoffice→docx")

    def test_битый_doc_не_роняет_приём(self):
        source = self.workdir / "обрывок.doc"
        source.write_bytes(os.urandom(2048))
        result = legacy.convert_doc(source)
        self.assertEqual(result.text, "")
        self.assertTrue(result.warnings, "битый файл обязан объясниться предупреждением")
        self.assertEqual(result.meta["source_format"], "doc")


@unittest.skipUnless(POSIX and pptx is not None, "нужен python-pptx и POSIX")
class LegacyPptTest(unittest.TestCase):
    """PPT → PPTX → конвертер презентаций из formats.office."""

    def setUp(self):
        self.workdir = Path(tempfile.mkdtemp(prefix="reportgen-test-"))
        self.addCleanup(shutil.rmtree, str(self.workdir), True)

    def test_слайды_из_ppt(self):
        payload = make_pptx(self.workdir)
        script = make_fake_soffice(self.workdir, payload=payload)
        source = self.workdir / "доклад.ppt"
        source.write_bytes(b"\xd0\xcf\x11\xe0 old powerpoint")
        with use_fake_soffice(script):
            result = legacy.convert_ppt(source)
        self.assertIn("## Слайд 1", result.text)
        self.assertIn("## Слайд 2", result.text)
        self.assertIn("Состояние линии", result.text)
        self.assertIn("Замена ODU", result.text)
        self.assertIn("Заметка к слайду 1", result.text)
        self.assertEqual(result.page_count, 2)
        self.assertEqual(result.meta["source_format"], "ppt")
        self.assertEqual(result.meta["converted_via"], "libreoffice→pptx")
        chunks = chunks_from_markdown(result.text, "reports/доклад", "reports")
        self.assertTrue(any(chunk.meta.get("page") == 2 for chunk in chunks),
                        "номер слайда должен доезжать до чанка")

    @unittest.skipUnless(LIBREOFFICE_IMPRESS,
                         "здешний LibreOffice без модуля Impress: .ppt не собрать")
    def test_настоящий_ppt_через_libreoffice(self):
        modern = make_pptx(self.workdir)
        legacy_path = self.workdir / "доклад.ppt"
        with legacy.convert_with_soffice(modern, "ppt") as produced:
            self.assertIsNotNone(produced.path, produced.warning)
            shutil.copyfile(str(produced.path), str(legacy_path))
        result = legacy.convert_ppt(legacy_path)
        self.assertIn("Состояние линии", result.text)
        self.assertIn("## Слайд 2", result.text)

    def test_без_модуля_office_понятное_предупреждение(self):
        """Модуль конвертера презентаций может не приехать в комплекте."""
        payload = make_pptx(self.workdir)
        script = make_fake_soffice(self.workdir, payload=payload)
        source = self.workdir / "доклад.ppt"
        source.write_bytes(b"\xd0\xcf\x11\xe0 old powerpoint")
        with use_fake_soffice(script), \
                mock.patch.dict(sys.modules, {"reportgen.ingest.formats.office": None}):
            result = legacy.convert_ppt(source)
        self.assertEqual(result.text, "")
        warning = " ".join(result.warnings)
        self.assertIn("office", warning)
        self.assertIn("python-pptx", warning)


@unittest.skipUnless(POSIX, "подставной soffice — исполняемый скрипт, только POSIX")
class LegacyOtherFormatsTest(unittest.TestCase):
    """XLS, PUB и запасной путь для OpenDocument."""

    def setUp(self):
        self.workdir = Path(tempfile.mkdtemp(prefix="reportgen-test-"))
        self.addCleanup(shutil.rmtree, str(self.workdir), True)
        self.log = self.workdir / "argv.log"

    @unittest.skipUnless(openpyxl is not None, "нужен openpyxl")
    def test_xls_через_libreoffice(self):
        payload = make_xlsx(self.workdir)
        script = make_fake_soffice(self.workdir, payload=payload, log=self.log)
        source = self.workdir / "бюджет.xls"
        source.write_bytes(b"\xd0\xcf\x11\xe0 old excel")
        with use_fake_soffice(script):
            result = legacy.convert_xls_soffice(source)
        self.assertIn("Бюджет канала", result.text)
        self.assertIn("| Параметр | Значение |", result.text)
        self.assertIn("52,4", result.text.replace(".", ","))
        self.assertEqual(result.meta["source_format"], "xls")
        self.assertEqual(option_value(read_log(self.log)[0], "--convert-to"), "xlsx")

    @unittest.skipUnless(pymupdf is not None and cyrillic_font(),
                         "нужен pymupdf и шрифт TTF с кириллицей")
    def test_pub_через_pdf(self):
        payload = make_pdf(self.workdir)
        script = make_fake_soffice(self.workdir, payload=payload, log=self.log)
        source = self.workdir / "буклет.pub"
        source.write_bytes(b"\xd0\xcf\x11\xe0 publisher")
        with use_fake_soffice(script):
            result = legacy.convert_pub(source)
        self.assertIn("Земная станция связи", result.text)
        self.assertEqual(result.page_count, 2)
        self.assertEqual(result.meta["converted_via"], "libreoffice→pdf")
        self.assertEqual(option_value(read_log(self.log)[0], "--convert-to"), "pdf")
        chunks = chunks_from_markdown(result.text, "literature/буклет", "literature")
        self.assertTrue(any(chunk.meta.get("page") == 2 for chunk in chunks))

    @unittest.skipUnless(openpyxl is not None, "нужен openpyxl")
    def test_opendocument_запасной_путь_выбирает_целевой_формат(self):
        payload = make_xlsx(self.workdir)
        script = make_fake_soffice(self.workdir, payload=payload, log=self.log)
        source = self.workdir / "смета.ods"
        source.write_bytes(b"PK\x03\x04 opendocument")
        with use_fake_soffice(script):
            result = legacy.convert_opendocument(source)
        self.assertEqual(option_value(read_log(self.log)[0], "--convert-to"), "xlsx")
        self.assertIn("| Параметр | Значение |", result.text)
        self.assertEqual(result.meta["source_format"], "ods")


@unittest.skipUnless(POSIX and docx is not None, "нужен python-docx и POSIX")
class ParallelConversionTest(unittest.TestCase):
    """Классическая ловушка: общий профиль ломает параллельные запуски."""

    def test_две_конвертации_в_потоках(self):
        workdir = Path(tempfile.mkdtemp(prefix="reportgen-test-"))
        self.addCleanup(shutil.rmtree, str(workdir), True)
        payload = make_docx(workdir)
        log = workdir / "argv.log"
        script = make_fake_soffice(workdir, payload=payload, sleep=0.3, log=log)
        sources = []
        for number in (1, 2):
            source = workdir / f"протокол-{number}.doc"
            source.write_bytes(b"\xd0\xcf\x11\xe0 old word document")
            sources.append(source)

        results = {}

        def run(source):
            results[source.name] = legacy.convert_doc(source)

        before = soffice_leftovers()
        with use_fake_soffice(script):
            threads = [threading.Thread(target=run, args=(item,)) for item in sources]
            for thread in threads:
                thread.start()
            for thread in threads:
                thread.join(timeout=60)

        self.assertEqual(len(results), 2)
        for result in results.values():
            self.assertEqual(result.warnings, [])
            self.assertIn("# Протокол измерений", result.text)
        profiles = [
            item
            for argv in read_log(log)
            for item in argv
            if item.startswith("-env:UserInstallation=")
        ]
        self.assertEqual(len(profiles), 2)
        self.assertEqual(len(set(profiles)), 2, "у каждого запуска обязан быть свой профиль")
        self.assertEqual(soffice_leftovers(), before)


# ------------------------------------------------------------------------ RTF ---

RTF_HEAD = (
    r"{\rtf1\ansi\ansicpg1251\deff0"
    "\n" r"{\fonttbl{\f0\froman Times New Roman;}}"
    "\n" r"{\stylesheet{\s1\fs32\b\sbasedon0 heading 1;}{\s2\fs28\b heading 2;}}"
    "\n"
)


def rtf_escape(text):
    """Кириллица в RTF пишется как \\'xx в кодовой странице документа."""
    out = []
    for character in text:
        if ord(character) < 128:
            out.append(character)
        else:
            out.append("\\'%02x" % character.encode("cp1251")[0])
    return "".join(out)


def make_rtf(body, head=RTF_HEAD):
    return head + body + "\n}"


class RtfTest(unittest.TestCase):
    """RTF через striprtf: заголовки, таблицы, списки, страницы, битые файлы."""

    def setUp(self):
        self.workdir = Path(tempfile.mkdtemp(prefix="reportgen-test-"))
        self.addCleanup(shutil.rmtree, str(self.workdir), True)

    def write(self, text, name="протокол.rtf", encoding="cp1251"):
        path = self.workdir / name
        path.write_bytes(text.encode(encoding))
        return path

    @unittest.skipUnless(striprtf is not None, "нужен striprtf")
    def test_заголовки_таблица_и_список(self):
        body = "\n".join([
            r"{\info{\title " + rtf_escape("Протокол измерений") + r"}{\author "
            + rtf_escape("Иванов") + r"}}",
            r"\pard\plain\s1\outlinelevel0 " + rtf_escape("Протокол измерений") + r"\par",
            r"\pard\plain " + rtf_escape(LONG_TEXT) + r"\par",
            r"\pard\plain\s2 " + rtf_escape("Допуски") + r"\par",
            r"\trowd\cellx3000\cellx6000",
            r"\intbl " + rtf_escape("Параметр") + r"\cell " + rtf_escape("Значение")
            + r"\cell\row",
            r"\trowd\cellx3000\cellx6000",
            r"\intbl EVM, %\cell 3,2\cell\row",
            r"\trowd\cellx3000\cellx6000",
            r"\intbl " + rtf_escape("Уровень приёма, дБм") + r"\cell "
            + rtf_escape("минус 42") + r"\cell\row",
            r"\pard{\listtext\f0 \'b7\tab}" + rtf_escape("Проверить юстировку") + r"\par",
            r"\pard{\listtext\f0 \'b7\tab}" + rtf_escape("Снять спектр") + r"\par",
        ])
        result = legacy.convert_rtf(self.write(make_rtf(body)))

        self.assertEqual(result.warnings, [])
        self.assertEqual(result.title, "Протокол измерений")
        self.assertEqual(result.meta["author"], "Иванов")
        self.assertEqual(result.meta["source_format"], "rtf")
        self.assertEqual(result.meta["converted_via"], "striprtf")
        self.assertIn("# Протокол измерений", result.text)
        self.assertIn("## Допуски", result.text)
        # Таблица допусков — целиком, одним блоком, с шапкой.
        self.assertIn(
            "| Параметр | Значение |\n| --- | --- |\n| EVM, % | 3,2 |\n"
            "| Уровень приёма, дБм | минус 42 |",
            result.text,
        )
        self.assertIn("- Проверить юстировку\n- Снять спектр", result.text)

    @unittest.skipUnless(striprtf is not None, "нужен striprtf")
    def test_разрывы_страниц_дают_номера(self):
        body = "\n".join([
            r"\pard\plain\s1\outlinelevel0 " + rtf_escape("Раздел первый") + r"\par",
            r"\pard " + rtf_escape(LONG_TEXT) + r"\par",
            r"\pard\page\s1\outlinelevel0 " + rtf_escape("Раздел второй") + r"\par",
            r"\pard " + rtf_escape("Измерения выполнены анализатором спектра. " + LONG_TEXT)
            + r"\par",
        ])
        result = legacy.convert_rtf(self.write(make_rtf(body)))
        self.assertEqual(result.page_count, 2)
        self.assertIn("<!-- page: 1 -->", result.text)
        self.assertIn("<!-- page: 2 -->", result.text)
        chunks = chunks_from_markdown(result.text, "reports/протокол", "reports")
        pages = {chunk.meta.get("page") for chunk in chunks}
        self.assertEqual(pages, {1, 2})

    @unittest.skipUnless(striprtf is not None, "нужен striprtf")
    def test_кириллица_без_экранирования(self):
        """Старые редакторы пишут кириллицу байтами cp1251, а не как \\'xx."""
        body = r"\pard\plain Затухание в дожде для диапазона Ku\par"
        result = legacy.convert_rtf(self.write(make_rtf(body)))
        self.assertIn("Затухание в дожде для диапазона Ku", result.text)
        self.assertEqual(result.meta["encoding"], "cp1251")

    @unittest.skipUnless(striprtf is not None, "нужен striprtf")
    def test_не_rtf_файл(self):
        path = self.workdir / "подделка.rtf"
        path.write_bytes("Это обычный текст, а не RTF".encode("utf-8"))
        result = legacy.convert_rtf(path)
        self.assertEqual(result.text, "")
        self.assertTrue(any("не похож на RTF" in item for item in result.warnings))

    @unittest.skipUnless(striprtf is not None, "нужен striprtf")
    def test_пустой_файл(self):
        path = self.workdir / "пусто.rtf"
        path.write_bytes(b"")
        result = legacy.convert_rtf(path)
        self.assertEqual(result.text, "")
        self.assertIn("файл пуст", result.warnings)

    @unittest.skipUnless(striprtf is not None, "нужен striprtf")
    def test_оборванный_файл_не_роняет_приём(self):
        """Файл скопирован не до конца: скобки не закрыты, хвост обрезан."""
        body = (
            r"\pard\plain\s1\outlinelevel0 " + rtf_escape("Начало документа") + r"\par"
            "\n" r"\trowd\cellx3000\intbl " + rtf_escape("Параметр") + r"\cell"
        )
        path = self.workdir / "обрывок.rtf"
        path.write_bytes((RTF_HEAD + body).encode("cp1251"))
        result = legacy.convert_rtf(path)
        self.assertIn("Начало документа", result.text)
        self.assertNotIn("@@RG-", result.text, "служебные метки не должны попадать в текст")

    @unittest.skipUnless(POSIX and docx is not None, "нужен python-docx и POSIX")
    def test_без_striprtf_идёт_через_libreoffice(self):
        payload = make_docx(self.workdir)
        script = make_fake_soffice(self.workdir, payload=payload, log=self.workdir / "argv.log")
        path = self.write(make_rtf(r"\pard\plain " + rtf_escape("Текст") + r"\par"))
        with mock.patch.dict(sys.modules, {"striprtf": None, "striprtf.striprtf": None}), \
                use_fake_soffice(script):
            result = legacy.convert_rtf(path)
        self.assertEqual(result.meta["source_format"], "rtf")
        self.assertEqual(result.meta["converted_via"], "libreoffice→docx")
        self.assertIn("# Протокол измерений", result.text)
        self.assertEqual(option_value(read_log(self.workdir / "argv.log")[0], "--convert-to"),
                         "docx")

    def test_без_striprtf_и_без_libreoffice_объясняет_чего_не_хватает(self):
        path = self.write(make_rtf(r"\pard\plain " + rtf_escape("Текст") + r"\par"))
        with mock.patch.dict(sys.modules, {"striprtf": None, "striprtf.striprtf": None}), \
                no_soffice():
            result = legacy.convert_rtf(path)
        self.assertEqual(result.text, "")
        warning = " ".join(result.warnings)
        self.assertIn("striprtf", warning)
        self.assertIn("LibreOffice", warning)


# ------------------------------------------------------------------- реестр ---

class RegistryTest(unittest.TestCase):
    """Как старые форматы видны диспетчеру и отчёту о поддержке."""

    def test_старые_форматы_зарегистрированы(self):
        for suffix, name in (
            (".doc", "doc-libreoffice"),
            (".dot", "doc-libreoffice"),
            (".wps", "doc-libreoffice"),
            (".ppt", "ppt-libreoffice"),
            (".pps", "ppt-libreoffice"),
            (".pub", "pub-libreoffice"),
            (".xlt", "xls-libreoffice"),
        ):
            spec = registry.find(suffix)
            self.assertIsNotNone(spec, f"формат {suffix} должен быть в реестре")
            self.assertEqual(spec.name, name, suffix)

    def test_требования_объявлены_с_подсказкой_для_windows(self):
        spec = registry.find(".doc")
        kinds = {item.name for item in spec.requires}
        self.assertIn("soffice", kinds)
        self.assertIn("docx", kinds)
        hint = next(item.hint for item in spec.requires if item.name == "soffice")
        self.assertIn("Windows", hint)
        self.assertIn("libreoffice", hint.lower())

    def test_требование_libreoffice_ищет_не_только_в_path(self):
        """В Windows soffice не в PATH, но формат всё равно поддержан."""
        spec = registry.find(".doc")
        requirement = next(item for item in spec.requires if item.name == "soffice")
        installed = r"C:\Program Files\LibreOffice\program\soffice.com"
        with mock.patch.object(legacy, "soffice_binary", lambda **kwargs: installed):
            self.assertTrue(requirement.is_available())
        with no_soffice():
            self.assertFalse(requirement.is_available())

    @unittest.skipUnless(xlrd is not None, "нужен xlrd, иначе .xls достаётся LibreOffice")
    def test_xls_приоритет_ниже_чем_у_xlrd(self):
        chosen = registry.find(".xls")
        self.assertEqual(chosen.name, "xls", "xlrd быстрее и не требует LibreOffice")
        ours = next(spec for spec in registry.all_specs() if spec.name == "xls-libreoffice")
        self.assertEqual(ours.priority, -10)
        self.assertLess(ours.priority, chosen.priority)

    def test_opendocument_запасной_путь_ниже_прямого_разбора(self):
        chosen = registry.find(".odt")
        ours = next(spec for spec in registry.all_specs()
                    if spec.name == "opendocument-libreoffice")
        self.assertEqual(ours.priority, -20)
        self.assertNotEqual(chosen.name, ours.name,
                            "прямой разбор OpenDocument не требует LibreOffice и всегда лучше")
        self.assertLess(ours.priority, chosen.priority)

    @unittest.skipUnless(striprtf is not None, "нужен striprtf")
    def test_rtf_striprtf_важнее_libreoffice(self):
        chosen = registry.find(".rtf")
        self.assertEqual(chosen.name, "rtf")
        fallback = next(spec for spec in registry.all_specs() if spec.name == "rtf-libreoffice")
        self.assertEqual(fallback.priority, -10)
        self.assertEqual(fallback.suffixes, (".rtf",))


class DispatchTest(unittest.TestCase):
    """Старые форматы через общую точку входа приёма — convert_file."""

    def setUp(self):
        self.workdir = Path(tempfile.mkdtemp(prefix="reportgen-test-"))
        self.addCleanup(shutil.rmtree, str(self.workdir), True)

    @unittest.skipUnless(striprtf is not None, "нужен striprtf")
    def test_rtf_разбирается_через_convert_file(self):
        path = self.workdir / "методика.rtf"
        body = (
            r"\pard\plain\s1\outlinelevel0 " + rtf_escape("Методика измерений") + r"\par"
            "\n" r"\pard " + rtf_escape(LONG_TEXT) + r"\par"
        )
        path.write_bytes(make_rtf(body).encode("cp1251"))
        result = convert_file(path)
        self.assertEqual(result.warnings, [])
        self.assertIn("# Методика измерений", result.text)
        self.assertEqual(result.meta["source_format"], "rtf")

    def test_без_libreoffice_convert_file_говорит_чего_не_хватает(self):
        path = self.workdir / "инструкция.doc"
        path.write_bytes(b"\xd0\xcf\x11\xe0 old word document")
        with no_soffice():
            result = convert_file(path)
        self.assertEqual(result.text, "")
        warning = " ".join(result.warnings)
        self.assertIn("doc-libreoffice", warning)
        self.assertIn("не хватает", warning)
        self.assertIn("soffice", warning)


if __name__ == "__main__":  # pragma: no cover — ручной запуск
    unittest.main()
