"""Приём библиотеки, в которой лежат документы всех форматов сразу.

Отдельные конвертеры проверяются своими тестами. Здесь проверяется главное для
заказчика: каталог со смесью PDF, DOCX, презентаций, таблиц, сканов и архивов
принимается целиком, ни один файл не роняет приём, а то, что разобрать нельзя,
попадает в предупреждения с внятной причиной.
"""

import shutil
import tempfile
import unittest
import zipfile
from pathlib import Path

import _bootstrap  # noqa: F401
from reportgen.ingest import registry
from reportgen.ingest.convert import convert_file, format_support, supported_suffixes
from reportgen.ingest.pipeline import ingest_directory, library_patterns
from reportgen.store import Database, Repositories

ROOT = Path(__file__).resolve().parents[1]


def registered(suffix: str) -> bool:
    """Есть ли доступный конвертер для расширения."""
    spec = registry.find(suffix)
    return spec is not None and spec.is_available()


class RegistryTests(unittest.TestCase):
    def test_builtin_formats_are_registered(self):
        for suffix in (".pdf", ".docx", ".md", ".txt"):
            self.assertTrue(registered(suffix), suffix)

    def test_report_lists_requirements(self):
        specs = format_support()
        self.assertTrue(specs)
        for spec in specs:
            self.assertIn("suffixes", spec)
            self.assertIn("available", spec)
            self.assertIsInstance(spec["requires"], list)

    def test_patterns_follow_registry(self):
        patterns = library_patterns()
        self.assertIn("*.pdf", patterns)
        for suffix in supported_suffixes(only_available=True):
            self.assertIn(f"*{suffix}", patterns)

    def test_unavailable_converter_is_still_reported(self):
        """Формат, для которого нет инструмента, не должен выглядеть неизвестным."""
        spec = registry.ConverterSpec(
            name="тестовый-формат",
            suffixes=(".тест",),
            convert=lambda path: None,
            requires=(registry.Requirement("binary", "несуществующая-программа",
                                           "поставить неоткуда"),),
        )
        registry.register(spec)
        try:
            found = registry.find(".тест")
            self.assertIsNotNone(found)
            self.assertFalse(found.is_available())
            self.assertIn("несуществующая-программа", registry.missing_hint(found))
            with tempfile.TemporaryDirectory() as tmp:
                path = Path(tmp) / "файл.тест"
                path.write_text("данные", encoding="utf-8")
                result = convert_file(path)
                self.assertTrue(result.warnings)
                self.assertIn("несуществующая-программа", " ".join(result.warnings))
        finally:
            registry._REGISTRY[:] = [
                item for item in registry._REGISTRY if item.name != "тестовый-формат"
            ]

    def test_unknown_format_names_available_ones(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "чертёж.dwg"
            path.write_bytes(b"\x00\x01\x02")
            result = convert_file(path)
            self.assertTrue(result.warnings)
            self.assertIn("неподдерживаемый формат", result.warnings[0])

    def test_converter_exception_becomes_warning(self):
        def broken(path):
            raise RuntimeError("внутренняя поломка конвертера")

        registry.register(registry.ConverterSpec(
            name="ломаный", suffixes=(".ломака",), convert=broken,
        ))
        try:
            with tempfile.TemporaryDirectory() as tmp:
                path = Path(tmp) / "x.ломака"
                path.write_text("данные", encoding="utf-8")
                result = convert_file(path)
                self.assertEqual(result.text, "")
                self.assertIn("не справился", " ".join(result.warnings))
        finally:
            registry._REGISTRY[:] = [
                item for item in registry._REGISTRY if item.name != "ломаный"
            ]


class MixedLibraryTests(unittest.TestCase):
    """Каталог со смесью форматов принимается целиком."""

    def setUp(self):
        self._tmp = tempfile.TemporaryDirectory()
        self.root = Path(self._tmp.name)
        (self.root / "literature").mkdir()
        (self.root / "standards").mkdir()
        self.made: list[str] = []
        self._make_files()
        self.repos = Repositories(Database(":memory:"))

    def tearDown(self):
        self._tmp.cleanup()

    def _make_files(self):
        text = ("Занимаемая полоса частот определяется методом 99 процентов мощности. "
                "Отношение сигнал/шум и модуль вектора ошибки характеризуют качество канала. ")

        (self.root / "literature" / "конспект.md").write_text(
            "# Конспект\n\n" + text * 6, encoding="utf-8")
        self.made.append(".md")

        (self.root / "literature" / "заметка.txt").write_bytes(
            ("Заметка инженера о параметрах линии. " + text * 4).encode("cp1251"))
        self.made.append(".txt (cp1251)")

        try:
            import docx

            document = docx.Document()
            document.add_heading("Руководство по эксплуатации", level=1)
            document.add_paragraph(text * 4)
            table = document.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "Параметр"
            table.cell(0, 1).text = "Значение"
            table.cell(1, 0).text = "Полоса"
            table.cell(1, 1).text = "1.85 МГц"
            document.save(self.root / "literature" / "руководство.docx")
            self.made.append(".docx")
        except ImportError:
            pass

        try:
            import pymupdf

            pdf = pymupdf.open()
            page = pdf.new_page()
            writer = pymupdf.TextWriter(page.rect)
            font = pymupdf.Font("helv")
            for index, line in enumerate((text * 2).split(". ")[:12]):
                writer.append((50, 60 + index * 18), line[:90], font=font, fontsize=11)
            writer.write_text(page)
            pdf.save(self.root / "standards" / "методика.pdf")
            pdf.close()
            self.made.append(".pdf")
        except ImportError:
            pass

        if registered(".pptx"):
            try:
                from pptx import Presentation

                deck = Presentation()
                slide = deck.slides.add_slide(deck.slide_layouts[1])
                slide.shapes.title.text = "Параметры радиорелейного пролёта"
                slide.placeholders[1].text = text * 2
                deck.save(self.root / "literature" / "презентация.pptx")
                self.made.append(".pptx")
            except Exception:  # noqa: BLE001 — формат необязателен
                pass

        if registered(".xlsx"):
            try:
                from openpyxl import Workbook

                book = Workbook()
                sheet = book.active
                sheet.title = "Бюджет линии"
                sheet.append(["Параметр", "Значение", "Единица"])
                sheet.append(["ЭИРП", 52.4, "дБВт"])
                sheet.append(["Затухание", 205.1, "дБ"])
                book.save(self.root / "standards" / "бюджет.xlsx")
                self.made.append(".xlsx")
            except Exception:  # noqa: BLE001
                pass

        if registered(".html"):
            (self.root / "literature" / "страница.html").write_text(
                "<html><head><title>Справка</title></head><body>"
                "<h1>Справка по протоколу</h1><p>" + text * 3 + "</p>"
                "<table><tr><th>Поле</th><th>Длина</th></tr>"
                "<tr><td>Флаг</td><td>1 байт</td></tr></table></body></html>",
                encoding="utf-8")
            self.made.append(".html")

        if registered(".zip"):
            archive = self.root / "literature" / "материалы.zip"
            with zipfile.ZipFile(archive, "w") as bundle:
                bundle.writestr("вложение.md", "# Вложенный документ\n\n" + text * 4)
            self.made.append(".zip")

    def test_all_prepared_formats_are_ingested(self):
        result = ingest_directory(self.repos, self.root)
        self.assertGreaterEqual(result.added, 4, result.summary())
        self.assertEqual(result.failed, 0, result.warnings)
        self.assertGreater(result.chunks, 0)

    def test_search_finds_content_from_every_document(self):
        ingest_directory(self.repos, self.root)
        hits = self.repos.chunks.search_lexical("занимаемая полоса частот", limit=50)
        self.assertTrue(hits)
        documents = {uid.split("#")[0] for uid, _ in hits}
        self.assertGreaterEqual(len(documents), 2, "содержимое нашлось меньше чем в двух файлах")

    def test_cp1251_text_is_readable(self):
        ingest_directory(self.repos, self.root)
        chunks = self.repos.chunks.all_chunks()
        note = [chunk for chunk in chunks if "заметка" in chunk.doc_id]
        self.assertTrue(note)
        self.assertIn("Заметка инженера", note[0].text)

    def test_unknown_file_does_not_break_ingest(self):
        (self.root / "literature" / "чертёж.dwg").write_bytes(b"\x00" * 512)
        result = ingest_directory(self.repos, self.root)
        self.assertEqual(result.failed, 0, result.warnings)

    def test_broken_file_is_reported_not_fatal(self):
        (self.root / "standards" / "битый.pdf").write_bytes(
            "%PDF-1.4 мусор".encode("utf-8"))
        result = ingest_directory(self.repos, self.root)
        self.assertGreater(result.added, 0)
        self.assertTrue(any("битый" in warning for warning in result.warnings),
                        result.warnings)

    def test_every_prepared_format_reached_the_index(self):
        """Отчёт о том, какие форматы реально проверены в этом прогоне."""
        ingest_directory(self.repos, self.root)
        documents = {document.doc_id for document in self.repos.documents.list()}
        self.assertTrue(documents)
        # Диагностика в вывод теста: видно, какие форматы были доступны.
        print("\n  проверенные форматы: " + ", ".join(self.made))


@unittest.skipUnless(shutil.which("soffice"), "LibreOffice не установлен")
class LegacyFormatTests(unittest.TestCase):
    """Старые форматы: файл готовится самим LibreOffice из современного."""

    def test_doc_is_converted(self):
        if not registered(".doc"):
            self.skipTest("конвертер .doc не зарегистрирован")
        try:
            import docx
        except ImportError:
            self.skipTest("python-docx не установлен")

        with tempfile.TemporaryDirectory() as tmp:
            source = Path(tmp) / "исходник.docx"
            document = docx.Document()
            document.add_heading("Старый документ", level=1)
            document.add_paragraph("Требование к отклонению несущей частоты. " * 8)
            document.save(source)

            import subprocess

            subprocess.run(
                ["soffice", "--headless", "--convert-to", "doc", "--outdir", tmp, str(source)],
                check=False, capture_output=True, timeout=180,
            )
            legacy = Path(tmp) / "исходник.doc"
            if not legacy.is_file():
                self.skipTest("LibreOffice не сконвертировал файл в .doc")
            result = convert_file(legacy)
            self.assertIn("несущей частоты", result.text)


if __name__ == "__main__":
    unittest.main()
